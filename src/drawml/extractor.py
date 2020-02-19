#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from math import sqrt, sin, cos, pi as PI
import os

#===============================================================================

import numpy as np

import pptx.shapes.connector
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

#===============================================================================

from ..parser import Parser

#===============================================================================

# Internal PPT units are EMUs (English Metric Units)

EMU_PER_CM  = 360000
EMU_PER_IN  = 914400

DOTS_PER_IN = 96

EMU_PER_DOT = EMU_PER_IN/DOTS_PER_IN

#===============================================================================

def cm_coords(x, y):
#===================
    return (x/EMU_PER_CM, y/EMU_PER_CM)

def ellipse_point(a, b, theta):
#==============================
    a_sin_theta = a*sin(theta)
    b_cos_theta = b*cos(theta)
    circle_radius = sqrt(a_sin_theta**2 + b_cos_theta**2)
    return (a*b_cos_theta/circle_radius, b*a_sin_theta/circle_radius)

#===============================================================================

class Transform(object):
    def __init__(self, shape, bbox=None):
        xfrm = shape.element.xfrm

        # From Section L.4.7.6 of ECMA-376 Part 1
        (Bx, By) = ((xfrm.chOff.x, xfrm.chOff.y)
                        if xfrm.chOff is not None else
                    (0, 0))
        (Dx, Dy) = ((xfrm.chExt.cx, xfrm.chExt.cy)
                        if xfrm.chExt is not None else
                    bbox)
        (Bx_, By_) = (xfrm.off.x, xfrm.off.y)
        (Dx_, Dy_) = (xfrm.ext.cx, xfrm.ext.cy)
        theta = xfrm.rot*PI/180.0
        Fx = -1 if xfrm.flipH else 1
        Fy = -1 if xfrm.flipV else 1
        T_st = np.array([[Dx_/Dx,      0, Bx_ - (Dx_/Dx)*Bx] if Dx != 0 else [1, 0, Bx_],
                         [     0, Dy_/Dy, By_ - (Dy_/Dy)*By] if Dy != 0 else [0, 1, By_],
                         [     0,      0,                 1]])
        U = np.array([[1, 0, -(Bx_ + Dx_/2.0)],
                      [0, 1, -(By_ + Dy_/2.0)],
                      [0, 0,                1]])
        R = np.array([[cos(theta), -sin(theta), 0],
                      [sin(theta),  cos(theta), 0],
                      [0,                    0, 1]])
        Flip = np.array([[Fx,  0, 0],
                         [ 0, Fy, 0],
                         [ 0,  0, 1]])
        T_rf = np.linalg.inv(U)@R@Flip@U
        self._T = T_rf@T_st

    def matrix(self):
        return self._T

#===============================================================================

class SlideToLayer(object):
    def __init__(self, extractor, slide, slide_number):
        self._extractor = extractor
        self._slide = slide
        self._slide_number = slide_number
        self._errors = []
        # Find `layer-id` text boxes so we have a valid ID **before** using
        # it when setting a shape's `path_id`.
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            if notes_text.startswith('.'):
                layer_directive = Parser.layer_directive(notes_text)
                if 'error' in layer_directive:
                    self._errors.append('Slide {}: invalid layer directive: {}'
                                        .format(slide_number, notes_text))
                    self._layer_id = 'layer-{:02d}'.format(slide_number)
                else:
                    self._layer_id = layer_directive.get('id')
                self._description = layer_directive.get('description', self._layer_id.capitalize())
                self._models = layer_directive.get('models', '')
                self._background_for = layer_directive.get('background-for', '')
                self._selectable = self._background_for == '' and not layer_directive.get('not-selectable')
                self._selected = layer_directive.get('selected', False)
                self._queryable_nodes = layer_directive.get('queryable-nodes', False)
                self._zoom = layer_directive.get('zoom', None)
            else:
                # still need to initialise properties...
                pass
        else:
            self._layer_id = 'layer-{:02d}'.format(slide_number)
            self._description = 'Layer {}'.format(slide_number)
            self._describes = ''
            self._background_for = ''
            self._selectable = False
            self._selected = False
            self._queryable_nodes = False
            self._zoom = None
        self._annotated_ids = {}
        self._map_features = []
        self._metadata = {}

    @property
    def settings(self):
        return self._extractor.settings

    @property
    def metadata(self):
        return self._metadata

    @property
    def description(self):
        return self._description

    @property
    def describes(self):
        return self._describes

    @property
    def background_for(self):
        return self._background_for

    @property
    def selected(self):
        return self._selected

    @property
    def selectable(self):
        return self._selectable

    @property
    def queryable_nodes(self):
        return self._queryable_nodes

    @property
    def zoom(self):
        return self._zoom

    @property
    def errors(self):
        return self._errors

    @property
    def map_features(self):
        return self._map_features

    @property
    def layer_id(self):
        return self._layer_id

    @property
    def slide_id(self):
        return self._slide.slide_id

    def process(self):
    #=================
        self.process_shape_list(self._slide.shapes)

    def save(self, filename=None):
    #=============================
        # Override in sub-class
        pass

    def process_group(self, group, *args):
    #=====================================
        self.process_shape_list(group.shapes, *args)

    def process_shape(self, shape, *args):
    #=====================================
        # Override in sub-class
        pass

    def process_shape_list(self, shapes, *args):
    #===========================================
        if not self._selectable:
            return
        for shape in shapes:
            shape.unique_id = '{}-{}'.format(self.slide_id, shape.shape_id)
            feature = None
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or shape.shape_type == MSO_SHAPE_TYPE.PICTURE
             or isinstance(shape, pptx.shapes.connector.Connector)):
                feature = self.process_shape(shape, *args)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.process_group(shape, *args)
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                pass
            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))

            if shape.name.startswith('#'):
                metadata = {
                    'layer': self.layer_id,
                    'annotation': shape.name
                }
                (annotated_id, properties) = Parser.annotation(shape.name)
                if annotated_id is not None:
                    if annotated_id in self._annotated_ids:
                        metadata['error'] = 'duplicate-id'
                        self._errors.append('Feature {} in slide {} has a duplicate feature id: {}'
                                            .format(shape.unique_id, self._slide_number, shape.name))
                    else:
                        self._annotated_ids[annotated_id] = shape.unique_id
                    if feature is not None:
                        label = properties.get('label', [None])[0]
                        models = properties.get('models', None)
                        if models is not None:
                            if label is None:
                                label = self.options.label_database.get_label(models[0][0])
                        if label is not None:
                            feature['properties']['label'] = label
                            metadata['label'] = label
                else:
                    metadata['error'] = 'syntax'
                    self._errors.append('Feature {} in slide {} has annotation syntax error: {}'
                                        .format(shape.unique_id, self._slide_number, shape.name))
                if feature is not None:
                    metadata['geometry'] = feature['geometry']['type']
                self._metadata[shape.unique_id] = metadata

#===============================================================================

class GeometryExtractor(object):
    def __init__(self, pptx, settings):
        self._pptx = Presentation(pptx)
        self._settings = settings
        self._slides = self._pptx.slides
        self._slide_size = [self._pptx.slide_width, self._pptx.slide_height]
        self._LayerMaker = SlideToLayer
        self._layers = {}

    def __len__(self):
        return len(self._slides)

    @property
    def layers(self):
        return self._layers

    @property
    def settings(self):
        return self._settings

    @property
    def slide_maker(self):
        return self._slide_maker

    @property
    def slide_size(self):
        return self._slide_size

    def bounds(self):
        return [0, 0, self._slide_size[0], self._slide_size[1]]

    def slide(self, slide_number):
        return self._slides[slide_number - 1]

    def slide_to_layer(self, slide_number, save_output=True):
        slide = self.slide(slide_number)
        if self._settings.debug_xml:
            xml = open(os.path.join(self._settings.output_dir, 'layer{:02d}.xml'.format(slide_number)), 'w')
            xml.write(slide.element.xml)
            xml.close()
        if self._LayerMaker is not None:
            layer = self._LayerMaker(self, slide, slide_number)
            layer.process()
            print('Slide {}, layer {}'.format(slide_number, layer.layer_id))
            if layer.layer_id in self._layers:
                raise KeyError('Duplicate layer id ({}) in slide {}'.format(layer.layer_id, slide_number))
            self._layers[layer.layer_id] = layer
            if save_output:
                layer.save()
            return layer

    def slides_to_layers(self, slide_range):
        if slide_range is None:
            slide_range = range(1, len(self._slides)+1)
        elif isinstance(slide_range, int):
            slide_range = [slide_range]
        for n in slide_range:
            self.slide_to_layer(n)

#===============================================================================
