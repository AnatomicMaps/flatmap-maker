#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from collections import OrderedDict

#===============================================================================

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tqdm import tqdm

#===============================================================================

from mapmaker.parser import Parser
from mapmaker.properties import JsonProperties

#===============================================================================

def shape_properties(shape):
    properties = {}
    if shape.name.startswith('.'):
        group_name = self.__current_group[-1]  # For error reporting
        properties.update(Parser.shape_markup(shape.name))
        if 'error' in properties:
            super().error('Shape in slide {}, group {}, has annotation syntax error: {}'
                          .format(self.__slide_number, group_name, shape.name))
        if 'warning' in properties:
            super().error('Warning, slide {}, group {}: {}'
                          .format(self.__slide_number, group_name, properties['warning']))
        for key in ['id', 'path']:
            if key in properties:
                if self.mapmaker.duplicate_id(properties[key]):
                   super().error('Shape in slide {}, group {}, has a duplicate id: {}'
                                 .format(self.__slide_number, group_name, shape.name))

class Slide(object):
    def __init__(self, slide, properties, options):
        self.__properties = properties
        self.__options = options
        self.__names_only = options.names_only
        self.__slide_id = slide.slide_id
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            self.__notes = notes_slide.notes_text_frame.text.rstrip()
        else:
            self.__notes = ''
        self.__properties_by_id = OrderedDict()
        self.__external_ids = []
        self.__current_group = ['SLIDE']
        self.process_shape_list_(slide.shapes, True)

    def process_shape_list_(self, shapes, outermost):
        if outermost and self.__options.verbose:
            progress_bar = tqdm(total=len(shapes),
                unit='shp', ncols=40,
                bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')
        for shape in shapes:
            unique_id = '{}#{}'.format(self.__slide_id, shape.shape_id)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                shape_type = 'G'
                self.__current_group.append(shape.name)
                self.process_shape_list_(shape.shapes, False)
                self.__current_group.pop()
            else:
                shape_type = 'S'
            print('{} {:10} {}'.format(shape_type, unique_id, shape.name))
##                self.__properties_by_id[unique_id] = self.__properties.get_properties(shape, self.__current_group[-1])
            if outermost and self.__options.verbose:
                progress_bar.update(1)
        if outermost and self.__options.verbose:
            progress_bar.close()

    def list(self):
        for id, properties in self.__properties_by_id.items():
            if properties and properties.get('shape-name', ' ')[0] in ['#', '.']:
                if self.__names_only:
                    print('{:8} {}'.format(id, properties['shape-name']))
                else:
                    print('{:8} {}'.format(id, properties))
        print('')

#===============================================================================

class Presentation(object):
    def __init__(self, powerpoint, properties, options):
        self.__pptx = pptx.Presentation(powerpoint)
        self.__slides_by_id = OrderedDict()
        self.__seen = []
        for slide_number, slide in enumerate(self.__pptx.slides):
            if options.debug_xml:
                xml = open(os.path.join(options.output_dir, 'layer{:02d}.xml'.format(slide_number)), 'w')
                xml.write(slide.element.xml)
                xml.close()
            self.__slides_by_id[slide.slide_id] = Slide(slide, properties, options)

    def list(self):
        for slide in self.__slides_by_id.values():
            slide.list()

#===============================================================================

if __name__ == '__main__':
    import configargparse
    import os, sys

    parser = configargparse.ArgumentParser() ## description='Modify annotations in Powerpoint slides.'

    parser.add_argument('-c', '--conf', is_config_file=True, help='configuration file containing arguments.')

#    parser.add_argument('-l', '--list', action='store_true',
#                        help='list shape annotations')

    parser.add_argument('--anatomical-map',
                        help='Excel spreadsheet file for mapping shape classes to anatomical entities')
    parser.add_argument('--properties',
                        help='JSON file specifying additional properties of shapes')

    parser.add_argument('-n', '--names-only', action='store_true',
                        help='Only list shape names')

    parser.add_argument('-v', '--verbose', action='store_true',
                        help='Show progress while running')
    parser.add_argument('-d', '--debug', dest='debug_xml', action='store_true',
                        help="save a slide's DrawML for debugging")

    # Options `--replace`
    # specify slide
    #         shape by id/class

    #parser.add_argument('-r', '--replace', nargs=2, action='append',
    #                    help='find and replace text using ')

    required = parser.add_argument_group('required arguments')

    required.add_argument('-o', '--output-dir', dest='map_base', metavar='OUTPUT_DIR', required=True,
                        help='base directory for generated flatmaps')
    required.add_argument('--id', dest='map_id', metavar='MAP_ID', required=True,
                        help='a unique identifier for the map')
    required.add_argument('--slides', dest='source', metavar='POWERPOINT', required=True,
                        help='Name of Powerpoint file to clean. The name of the resulting cleaned'
                        ' Powerpoint has `_cleaned` added to the source name.')

    args, unknown_args = parser.parse_known_args()

    args.label_database = 'labels.sqlite'
    json_properties = JsonProperties(args)

    args.output_dir = os.path.join(args.map_base, args.map_id)

    presentation = Presentation(args.source, json_properties, args)
    ##presentation.list()
