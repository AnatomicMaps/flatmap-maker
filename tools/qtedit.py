#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

#===============================================================================

class IdentifiedObject():
    def __init__(self, id):
        self._id = id

    @property
    def id(self):
        return self._id

#===============================================================================

class ShapeContainer(object):
    def __init__(self):
        self._shape_ids = []

    @property
    def shape_ids(self):
        return self._shape_ids

    def add_shape(self, shape):
        self._shape_ids.append(shape.id)

#===============================================================================

class Shape(IdentifiedObject):
    def __init__(self, slide_id, shape):
        IdentifiedObject.__init__(self, '{}-{}'.format(slide_id, shape.shape_id))
        self._shape = shape

    def __str__(self):
        return 'SHAPE: {!s:8} {}'.format(self.id, self._shape.name)

#===============================================================================

class Group(Shape, ShapeContainer):
    def __init__(self, slide_id, shape):
        Shape.__init__(self, slide_id, shape)
        ShapeContainer.__init__(self)

    def __str__(self):
        return 'GROUP: {!s:8} {} {}'.format(self.id, self.shape_ids, self._shape.name)

#===============================================================================

class Slide(IdentifiedObject, ShapeContainer):
    def __init__(self, slide):
        IdentifiedObject.__init__(self, slide.slide_id)
        ShapeContainer.__init__(self)
        self._notes = (slide.notes_slide.notes_text_frame.text.rstrip() if slide.has_notes_slide
                  else '')
        self._shapes_by_id = {}
        self.add_shapes_(self, slide.shapes)

    def __str__(self):
        return 'SLIDE: {!s:8} {} {}'.format(slide.id, slide.shape_ids, self._notes)

    def add_shapes_(self, group, shapes):
        for shape in shapes:
            new_shape = (Group(self._id, shape) if shape.shape_type == MSO_SHAPE_TYPE.GROUP
                    else Shape(self._id, shape))
            self._shapes_by_id[new_shape.id] = new_shape
            group.add_shape(new_shape)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.add_shapes_(new_shape, shape.shapes)

    def shape(self, id):
        return self._shapes_by_id.get(id)

#===============================================================================

class Presentation(object):
    def __init__(self, powerpoint):
        self._pptx = pptx.Presentation(powerpoint)
        self._slides = []
        for slide in self._pptx.slides:
            self._slides.append(Slide(slide))

#===============================================================================

if __name__ == '__main__':
    import argparse
    import os, sys

    parser = argparse.ArgumentParser(description='Modify annotations in Powerpoint slides.')

    parser.add_argument('-l', '--list', action='store_true',
                        help='list annotations')

    parser.add_argument('powerpoint', help='Powerpoint file')

    args = parser.parse_args()

    presentation = Presentation(args.powerpoint)

    slide = presentation._slides[0]
    print(slide)
    for shape in presentation._slides[0]._shapes_by_id.values():
        print(shape)

#===============================================================================
