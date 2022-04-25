#===============================================================================
#
#  Cell Diagramming Language
#
#  Copyright (c) 2018 - 2022  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from dataclasses import dataclass, field
from typing import Dict, List, Optional, Tuple

#===============================================================================

from lxml import etree              # type: ignore
import numpy as np
import shapely.geometry             # type: ignore
import shapely.ops                  # type: ignore

from pptx import Presentation       # type: ignore
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE, MSO_THEME_COLOR   # type: ignore
from pptx.enum.shapes import MSO_SHAPE_TYPE              # type: ignore
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN          # type: ignore
import pptx.shapes.connector                             # type: ignore

#===============================================================================

from mapmaker.geometry import Transform
from mapmaker.utils import FilePath, ProgressBar, log

from .. import WORLD_METRES_PER_EMU
from ..powerpoint.colour import ColourMap, Theme
from ..powerpoint.presets import DML
from ..powerpoint.transform import DrawMLTransform
from ..powerpoint.utils import get_shape_geometry

#===============================================================================

PPTX_NAMESPACE = {
    'p': "http://schemas.openxmlformats.org/presentationml/2006/main",
    'a': "http://schemas.openxmlformats.org/drawingml/2006/main",
    'r': "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
}

def pptx_resolve(qname):
    parts = qname.split(':', 1)
    if len(parts) == 2 and parts[0] in PPTX_NAMESPACE:
        return f'{{{PPTX_NAMESPACE[parts[0]]}}}{parts[1]}'
    return qname

#===============================================================================

def text_alignment(shape):
#=========================
    para = shape.text_frame.paragraphs[0].alignment
    vertical = shape.text_frame.vertical_anchor
    return ('left' if para in [PP_ALIGN.LEFT, PP_ALIGN.DISTRIBUTE, PP_ALIGN.JUSTIFY, PP_ALIGN.JUSTIFY_LOW] else
            'right' if para == PP_ALIGN.RIGHT else
            'centre',
            'top' if vertical == MSO_ANCHOR.TOP else
            'bottom' if vertical == MSO_ANCHOR.BOTTOM else
            'middle')

def text_content(shape):
#=======================
    text = shape.text.replace('\n', ' ').replace('\xA0', ' ').replace('\v', ' ').strip() # Newline, non-breaking space, vertical-tab
    return text if text not in ['', '.'] else ''

#===============================================================================

@dataclass
class Shape:
    type: str
    id: int
    geometry: shapely.geometry.base.BaseGeometry
    properties: Dict[str, str] = field(default_factory=dict)

    @property
    def label(self):
        return self.properties.get('label', '')

#===============================================================================

class Powerpoint():
    def __init__(self, source_href):
        ppt_bytes = FilePath(source_href).get_BytesIO()
        pptx = Presentation(ppt_bytes)
        theme = Theme(ppt_bytes)
        slides = pptx.slides
        (width, height) = (pptx.slide_width, pptx.slide_height)
        self.__transform = Transform([[WORLD_METRES_PER_EMU,                     0, 0],
                                      [                    0, -WORLD_METRES_PER_EMU, 0],
                                      [                    0,                     0, 1]])@np.array([[1, 0, -width/2.0],
                                                                                                    [0, 1, -height/2.0],
                                                                                                    [0, 0,         1.0]])
        self.__slide = Slide(slides[0], theme, self.__transform)
        top_left = self.__transform.transform_point((0, 0))
        bottom_right = self.__transform.transform_point((width, height))
        # southwest and northeast corners
        self.geometry = shapely.geometry.box(top_left[0], bottom_right[1], bottom_right[0], top_left[1])

    @property
    def transform(self):
        return self.__transform

    def process(self):
    #=================
        return self.__slide.process()

#===============================================================================

class Slide():
    def __init__(self, slide, theme, transform):
        self.__colour_map = ColourMap(theme, slide)
        self.__slide = slide
        self.__transform = transform

    def process(self):
    #=================
        return self.__process_shape_list(self.__slide.shapes, self.__transform, show_progress=True)

    def __get_colour(self, shape):
    #=============================
        colour = None
        alpha = 1.0
        if not isinstance(shape, pptx.shapes.connector.Connector):
            if shape.fill.type == MSO_FILL_TYPE.SOLID:
                colour = self.__colour_map.lookup(shape.fill.fore_color)
                alpha = shape.fill.fore_color.alpha
            elif shape.fill.type == MSO_FILL_TYPE.GRADIENT:
                log.warning(f'{shape.name}: gradient fill ignored')
            elif shape.fill.type == MSO_FILL_TYPE.GROUP:
                # WIP Need to get group's fill
                log.warning(f'{shape.name}: group fill ignored')
            elif shape.fill.type is not None and shape.fill.type != MSO_FILL_TYPE.BACKGROUND:
                log.warning(f'{shape.name}: unsupported fill type: {shape.fill.type}')
        elif shape.line.fill.type == MSO_FILL_TYPE.SOLID:
            colour = self.__colour_map.lookup(shape.line.color)
            alpha = shape.line.fill.fore_color.alpha
        elif shape.line.fill.type is None:
            # Check for a fill colour in the <style> block
            xml = etree.fromstring(shape.element.xml)
            if (scheme_colour := xml.find('.//p:style/a:fillRef/a:schemeClr',
                                            namespaces=PPTX_NAMESPACE)) is not None:
                colour = self.__colour_map.scheme_colour(scheme_colour.attrib['val'])
        elif shape.line.fill.type != MSO_FILL_TYPE.BACKGROUND:
            log.warning(f'{shape.name}: unsupported line fill type: {shape.line.fill.type}')
        return colour

    def __process_group(self, group, transform):
    #===========================================
        shapes = self.__process_shape_list(group.shapes, transform@DrawMLTransform(group))
        if len(shapes) < 2 or shapes[0].type != 'feature':
            return shapes
        colour = shapes[0].properties['colour']
        label = shapes[0].label
        alignment = shapes[0].properties.get('align')
        geometry = [shapes[0].geometry]
        for shape in shapes[1:]:
            if shape.type != 'feature' or colour != shape.properties['colour']:
                return shapes
            if label == '':
                label = shape.label
                alignment = shape.properties.get('align')
            elif shape.label != '':
                return shapes
            geometry.append(shape.geometry)
        if label == '':
            return shapes
        return [Shape('feature', group.shape_id,
                      shapely.ops.unary_union(geometry), {
                        'colour': colour,
                        'label': label,
                        'shape-name': group.name,
                        'text-align': alignment
                       })]

    def __process_shape_list(self, shape_list, transform, show_progress=False):
    #==========================================================================
        progress_bar = ProgressBar(show=show_progress,
            total=len(shape_list),
            unit='shp', ncols=40,
            bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')
        shapes = []
        for shape in shape_list:
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
             or isinstance(shape, pptx.shapes.connector.Connector)):
                geometry = get_shape_geometry(shape, transform)
                if geometry is not None:
                    shape_properties = {
                        'shape-name': shape.name,
                        'colour': self.__get_colour(shape)
                    }
                    shape_xml = etree.fromstring(shape.element.xml)
                    shape_links = set()
                    for link_ref in shape_xml.findall('.//a:hlinkClick',
                                                    namespaces=PPTX_NAMESPACE):
                        r_id = link_ref.attrib[pptx_resolve('r:id')]
                        if (r_id in shape.part.rels
                         and shape.part.rels[r_id].reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink'):
                            shape_links.add(shape.part.rels[r_id].target_ref)
                    if len(shape_links):
                        shape_properties['hyperlinks'] = list(shape_links)
                    if isinstance(shape, pptx.shapes.connector.Connector):
                        shape_type = 'connector'
                        xml = etree.fromstring(shape.element.xml)
                        if (connection := shape_xml.find('.//p:x`/p:cNvCxnSpPr',
                                                        namespaces=PPTX_NAMESPACE)) is not None:
                            for c in connection.getchildren():
                                if c.tag == DML('stCxn'):
                                    shape_properties['connection-start'] = int(c.attrib['id'])
                                elif c.tag == DML('endCxn'):
                                    shape_properties['connection-end'] = int(c.attrib['id'])
                        line_style = 'solid'
                        head_end = 'none'
                        tail_end = 'none'
                        if (line_props := shape_xml.find('.//p:spPr/a:ln',
                                                        namespaces=PPTX_NAMESPACE)) is not None:
                            for prop in line_props.getchildren():
                                if prop.tag == DML('prstDash'):
                                    line_style = prop.attrib['val']
                                elif prop.tag == DML('headEnd'):
                                    head_end = prop.attrib['type']
                                elif prop.tag == DML('tailEnd'):
                                    tail_end = prop.attrib['type']
                        shape_properties['line-style'] = line_style
                        shape_properties['head-end'] = head_end
                        shape_properties['tail-end'] = tail_end
                    else:
                        shape_type = 'feature'
                        label = text_content(shape)
                        if label != '':
                            shape_properties['label'] = label
                            shape_properties['align'] = text_alignment(shape)
                    shapes.append(Shape(shape_type, shape.shape_id, geometry, shape_properties))
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                shapes.extend(self.__process_group(shape, transform))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                log.warning('Image "{}" {} not processed...'.format(shape.name, str(shape.shape_type)))
            else:
                log.warning('Shape "{}" {} not processed...'.format(shape.name, str(shape.shape_type)))
            progress_bar.update(1)

        progress_bar.close()
        return shapes

#===============================================================================
