#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019 - 2023  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from io import StringIO
from math import sqrt
from typing import Any, Optional

#===============================================================================

from lxml import etree
import numpy as np

import svgelements
import svgwrite.gradients
from svgwrite import Drawing as SvgDrawing
from svgwrite.base import BaseElement as SvgElement
from svgwrite.base import Desc as SvgDescription
from svgwrite.container import Defs as SvgDefinitions
from svgwrite.container import Group as SvgGroup
from svgwrite.container import Hyperlink as SvgHyperlink
from svgwrite.image import Image as SvgImage
from svgwrite.path import Path as SvgPath
from svgwrite.text import Text as SvgText

#===============================================================================

from pptx.dml.fill import FillFormat
from pptx.enum.dml import MSO_FILL_TYPE, MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.util import Length

from pptx.shapes.autoshape import Shape as PptxShape
from pptx.shapes.connector import Connector as PptxConnector
from pptx.shapes.group import GroupShape as PptxGroupShape

#===============================================================================

from mapmaker.geometry import Transform
from mapmaker.knowledgebase.celldl import CD_CLASS, FC_CLASS
from mapmaker.knowledgebase.celldl import CellDLGraph
from mapmaker.settings import settings
from mapmaker.sources import EMU_PER_METRE, MapBounds, WORLD_METRES_PER_PIXEL, POINTS_PER_PIXEL
from mapmaker.sources.shape import Shape, SHAPE_TYPE
from mapmaker.utils import log, TreeList
from mapmaker.utils.svg import css_class, name_from_id, svg_id

from .colour import ColourPair, ColourMap
from .presets import DRAWINGML, PPTX_NAMESPACE, pptx_resolve, pptx_uri
from .powerpoint import Powerpoint, Slide

#===============================================================================

# Minimum width for a stroked path in points
MIN_STROKE_WIDTH = 0.5

TEXT_MARGINS = (6, 0)   # pixels

#===============================================================================

def text_alignment(shape):
#=========================
    para = shape.text_frame.paragraphs[0].alignment
    vertical = shape.text_frame.vertical_anchor
    return ('left' if para in [PP_ALIGN.LEFT, PP_ALIGN.DISTRIBUTE, PP_ALIGN.JUSTIFY, PP_ALIGN.JUSTIFY_LOW] else  # type: ignore
            'right' if para == PP_ALIGN.RIGHT else              # type: ignore
            'centre',
            'top' if vertical == MSO_ANCHOR.TOP else            # type: ignore
            'bottom' if vertical == MSO_ANCHOR.BOTTOM else      # type: ignore
            'top')

def text_content(shape):
#=======================
    text = shape.text.replace('\n', ' ').replace('\xA0', ' ').replace('\v', ' ').strip() # Newline, non-breaking space, vertical-tab
    return text if text not in ['', '.'] else None

#===============================================================================

def points_to_pixels(pts):
#======================
    return pts/POINTS_PER_PIXEL

#===============================================================================

## NB. Adobe Illustrator 2020 doesn't appear to support marker definitions in SVG

def __scale(t):
#==============
    try:
        return str(POINTS_PER_PIXEL*float(t))
    except ValueError:
        return t

def __scale_to_world(s: str) -> str:
#===================================
    return ' '.join([__scale(t) for t in s.split()])

ARROW_MARKERS = {
    'triangle-head': __scale_to_world('M 10 0 L 0 5 L 10 10 z'),
    'triangle-tail': __scale_to_world('M 0 0 L 10 5 L 0 10 z')
}

def add_marker_definitions(drawing: SvgDrawing):
#===============================================
    # arrowhead markers (see https://developer.mozilla.org/en-US/docs/Web/SVG/Element/marker)
    # 18 Jan 2023: markers appear in Chrome with black fill; no markers in Firefox
    for id, path in ARROW_MARKERS.items():
        marker = drawing.marker(id=svg_id(id),
                                viewBox=__scale_to_world('0 0 10 10'),
                                refX=__scale_to_world('5'),
                                refY=__scale_to_world('5'),
                                markerUnits='userSpaceOnUse',
                                markerWidth=__scale_to_world('6'),
                                markerHeight=__scale_to_world('6'),
                                orient='auto')
        marker.add(drawing.path(d=path))   ## , fill='context-stroke' is not supported by svgwrite
        drawing.defs.add(marker)

def marker_id(marker_def, end):
#==============================
    marker_type = marker_def.get('type', 'none')
    return ('#{}-{}'.format(marker_type, end)
            if marker_type != 'none'
            else None)

## only add definitions for markers that have been used??
## put the above into a class??

#===============================================================================

# Don't set a path id for default shape names

EXCLUDED_NAME_PREFIXES = [
    'Freeform',
    'Group',
    'Oval',
    'Star',
]

# Markup that has been deprecated

EXCLUDED_NAME_MARKUP = [
    '.siblings',
]

# Check to see if we have a valid name and encode it as an id

def valid_markup(name):
#======================
    if name not in EXCLUDED_NAME_MARKUP:
        for prefix in EXCLUDED_NAME_PREFIXES:
            if name.startswith(prefix):
                return False
        return True
    return False

def add_markup(element, markup):
#===============================
    if valid_markup(markup):
        element.set_desc(title=markup)

def add_class(xml, cls):
#=======================
    if 'class' not in xml.attribs:
        xml.attribs['class'] = cls
    else:
        xml.attribs['class'] += f' {cls}'

#===============================================================================

class Gradient(object):
    def __init__(self, definitions: SvgDefinitions, id: int, pptx_shape, colour_map: ColourMap):
        self.__id = 'gradient-{}'.format(id)
        fill = pptx_shape.fill
        gradient = None
        if fill._fill._gradFill.path is None:
            gradient = svgwrite.gradients.LinearGradient(id=svg_id(self.__id))
            rotation = fill.gradient_angle
            if ('rotWithShape' in fill._fill._gradFill.attrib
             and fill._fill._gradFill.attrib['rotWithShape'] == '1'):
                rotation += pptx_shape.rotation
            if rotation != 0:
                gradient.rotate(rotation % 360, (0.5, 0.5))

        elif fill._fill._gradFill.path.attrib['path'] == 'circle':
                fill_to = fill._fill._gradFill.path.find(DRAWINGML('fillToRect')).attrib
                tileRect = fill._fill._gradFill.find(DRAWINGML('tileRect'))
                tile = tileRect.attrib if tileRect is not None else {}
                cx = (float(fill_to['l']) if 'l' in fill_to else float(fill_to['r']) + float(tile.get('l', 0.0)))/100000.0
                cy = (float(fill_to['t']) if 't' in fill_to else float(fill_to['b']) + float(tile.get('t', 0.0)))/100000.0
                sx = (float(fill_to['r']) if 'r' in fill_to else float(fill_to['l']) + float(tile.get('r', 0.0)))/100000.0
                sy = (float(fill_to['b']) if 'b' in fill_to else float(fill_to['t']) + float(tile.get('b', 0.0)))/100000.0
                if pptx_shape.width > pptx_shape.height:
                    scale_x = pptx_shape.height/pptx_shape.width
                    scale_y = 1.0
                elif pptx_shape.width < pptx_shape.height:
                    scale_x = 1.0
                    scale_y = pptx_shape.width/pptx_shape.height
                else:
                    scale_x = 1.0
                    scale_y = 1.0
                if len(tile) == 0:
                    radius = 1.0
                    if (cx, cy) != (0.5, 0.5) and (sx, sy) != (0.5, 0.5):
                        print('Preset radial gradient for shape:', pptx_shape.name)
                else:
                    radius = sqrt(((cx-sx)/scale_x)**2 + ((cy-sy)/scale_y)**2)
                gradient = svgwrite.gradients.RadialGradient((cx/scale_x, cy/scale_y), radius, id=svg_id(self.__id))
                if pptx_shape.rotation != 0:
                    gradient.rotate(pptx_shape.rotation, (0.5, 0.5))
                if pptx_shape.width != pptx_shape.height:
                    gradient.scale(scale_x, scale_y)

        elif fill._fill._gradFill.path.attrib['path'] == 'rect':
            print('Rect fill ignored for', pptx_shape.name)
            return

        if gradient is not None:
            for stop in sorted(fill.gradient_stops, key=lambda stop: stop.position):
                gradient.add_stop_color(offset=stop.position,
                    color=colour_map.lookup(stop.color),
                    opacity=stop.color.alpha)
            definitions.add(gradient)
        else:
            print('UNKNOWN FILL: {}\n'.format(pptx_shape.name), fill._fill._element.xml)

    @property
    def url(self):
        return 'url(#{})'.format(self.__id)

#===============================================================================

class SvgFromSlide:
    def __init__(self, slide_id: str, slide: Slide, drawing: SvgDrawing,
                 transform_matrix: svgelements.Matrix, celldl: Optional[CellDLGraph]=None):
        self.__slide_id = slide_id
        self.__slide = slide
        self.__drawing = drawing
        self.__matrix = transform_matrix
        self.__celldl = celldl
        self.__colour_map = slide.colour_map
        self.__gradient_id = 0

    def process_slide_svgs(self):
    #============================
        slide_group = SvgGroup(id=svg_id(self.__slide_id), class_=css_class(CD_CLASS.LAYER))
        slide_group.set_desc(title=name_from_id(self.__slide_id))
        self.__process_shape_list(self.__slide.shapes, slide_group)
        self.__drawing.add(slide_group)

    def __process_group(self, group: TreeList, svg_parent: SvgElement):
    #==================================================================
        if group[0].type != SHAPE_TYPE.GROUP:
            raise TypeError(f'Invalid shape treelist: index 0 shape type ({group[0].type}) != SHAPE_TYPE.GROUP')
        svg_group = SvgGroup(id=svg_id(group[0].id))
        pptx_group = group[0].properties['pptx-shape']
        self.__process_shape_list(group, svg_group, group_colour=self.__get_colour(pptx_group))
        if len(svg_group.elements):
            svg_parent.add(svg_group)
            add_markup(svg_group, pptx_group.name)

    def __process_shape_list(self, shapes: TreeList, svg_parent: SvgElement,
                                   group_colour: Optional[ColourPair]=None):
    #=======================================================================
        for shape in shapes[1:]:
            if isinstance(shape, TreeList):
                self.__process_group(shape, svg_parent)
            elif shape.type in [SHAPE_TYPE.CONNECTION, SHAPE_TYPE.FEATURE]:
                self.__process_shape(shape, svg_parent, group_colour=group_colour)
            else:
                raise TypeError(f'Unexpected shape type: {shape.type}')

    def __process_shape(self, shape: Shape, svg_parent: SvgElement,
                              group_colour: Optional[ColourPair]=None):
    #==================================================================
        pptx_shape = shape.properties.pop('pptx-shape')
        svg_shape = shape.properties.pop('svg-element')
        svg_kind = shape.properties.pop('svg-kind')
        if svg_kind == 'group':
            svg_group = SvgGroup(id=svg_id(shape.id))
            for n, svg in enumerate(svg_shape):
                self.__process_svg_shape(svg, 'path', svg_group, group_colour,
                                         f'{shape.id}/{n}', shape.geometry.bounds, shape.properties,
                                         pptx_shape)
            if len(svg_group.elements):
                svg_parent.add(svg_group)
        else:
            self.__process_svg_shape(svg_shape, svg_kind, svg_parent, group_colour,
                                     shape.id, shape.geometry.bounds, shape.properties,
                                     pptx_shape)
        if self.__celldl is not None:
            self.__celldl.add_metadata(shape)

    def __process_svg_shape(self, svg_shape, svg_kind, svg_parent, group_colour,
                            shape_id, bounds, properties, pptx_shape):
    #=================================================================
        transformed_svg = svg_shape * self.__matrix
        transformed_svg.reify()

        if svg_kind == 'path':
            svg_element = SvgPath(str(transformed_svg),
                class_='non-scaling-stroke',
                fill='none',
                id=svg_id(shape_id))
        elif svg_kind == 'image':
            image_href = transformed_svg.values['data-image-href']
            bbox = transformed_svg.bbox()
            svg_element = SvgImage(image_href,
                insert=(bbox[0], bbox[1]),
                size=(bbox[2]-bbox[0], bbox[3]-bbox[1]),
                id=svg_id(shape_id))
        else:
            log.error(f'Unknown kind of SVG element for shape: {svg_kind}')
            return




        if not properties.get('exclude', False):   # Will this handle shape filtering??
                                                   # this does remove provenance stars but we need to capture links as FTU metadata
            celldl_class = properties.get('cd-class')
            if celldl_class is not None:
                svg_element.attribs['class'] = css_class(celldl_class)

            if celldl_class == CD_CLASS.CONNECTION:
                if 'type' in pptx_shape.line.headEnd or 'type' in pptx_shape.line.tailEnd:          # type: ignore
                    svg_element.set_markers((marker_id(pptx_shape.line.headEnd, 'head'),            # type: ignore
                                             None, marker_id(pptx_shape.line.tailEnd, 'tail')       # type: ignore
                                           ))
                svg_element.attribs.update(self.__get_stroke(pptx_shape))

                svg_parent.add(svg_element)

            elif svg_kind == 'image':
                svg_parent.add(svg_element)

            else:
                svg_element.attribs.update(self.__get_fill(pptx_shape, group_colour))
                svg_element.attribs.update(self.__get_stroke(pptx_shape))

                label = text_content(pptx_shape)    ### shape.label
                exclude_text = properties.get('fc-class') != FC_CLASS.SYSTEM
                svg_text = None
                if not exclude_text and label is not None:
                    svg_text = self.__draw_shape_label(pptx_shape, label, bounds)    # type: ignore
                if (hyperlink := self.__get_link(pptx_shape)) is not None:
                    if label is None:
                        label = hyperlink
                    link_element = SvgHyperlink(href=hyperlink)
                    link_element.add(svg_element)
                    if svg_text is not None:
                        link_element.add(svg_text)
                    svg_parent.add(link_element)
                else:
                    svg_parent.add(svg_element)
                    if svg_text is not None:
                        svg_parent.add(svg_text)
                if label is not None:
                    add_markup(svg_element, label)  # Set's <title>
                    if svg_text is not None:
                        add_markup(svg_text, label)   ## shape.label
                else:
                    add_markup(svg_element, pptx_shape.name)

    def __draw_shape_label(self, pptx_shape: PptxShape, label: str, bbox: MapBounds) -> SvgElement:
    #==============================================================================================
        shape_pos = (bbox[0], bbox[3])
        shape_size = (bbox[2]-bbox[0], bbox[3]-bbox[1])

        style = {}
        font = pptx_shape.text_frame.paragraphs[0].runs[0].font     # type: ignore
        font_size = round(font.size/EMU_PER_METRE)                  # type: ignore

        style['font-family'] = font.name if font.name is not None else 'Calibri'
        style['font-size'] = f'{font_size}px'
        style['font-weight'] = 700 if font.bold else 400
        if font.italic:
            style['font-size'] = 'italic'
        if font.color.type is not None:                             # type: ignore
            style['fill'] = self.__colour_map.lookup(font.color)
            if font.color.alpha != 1.0:                             # type: ignore
                style['fill-opacity'] = font.color.alpha            # type: ignore

        svg_text = SvgText(label)   ## text_run.text
        svg_text.attribs['style'] = ' '.join([f'{name}: {value};' for name, value in style.items()])

        svg_text.attribs['x'] =  shape_pos[0]
        svg_text.attribs['y'] = -shape_pos[1]
        (halign, valign) = text_alignment(pptx_shape)

        if pptx_shape.rotation != 0:
            rotation = f'rotate({-pptx_shape.rotation}, {shape_pos[0]+shape_size[0]/2.0}, {shape_pos[1]-shape_size[1]/2.0})'
        else:
            rotation = ''


        if halign == 'right':
            svg_text.attribs['text-anchor'] = 'end'
            svg_text.attribs['x'] += shape_size[0] - TEXT_MARGINS[0]*WORLD_METRES_PER_PIXEL
        elif halign == 'centre':
            svg_text.attribs['text-anchor'] = 'middle'
            svg_text.attribs['x'] += shape_size[0]/2
        else:   # Default to 'left'
            svg_text.attribs['text-anchor'] = 'start'
            svg_text.attribs['x'] += TEXT_MARGINS[0]*WORLD_METRES_PER_PIXEL
        if valign == 'bottom':
            svg_text.attribs['dominant-baseline'] = 'auto'
            svg_text.attribs['y'] += shape_size[1]
        elif valign == 'middle':
            svg_text.attribs['dominant-baseline'] = 'middle'
            svg_text.attribs['y'] += shape_size[1]/2
        else:   # Default to 'top'
            svg_text.attribs['dominant-baseline'] = 'auto'
            svg_text.attribs['y'] += font_size + TEXT_MARGINS[1]*WORLD_METRES_PER_PIXEL

        text_group = SvgGroup(transform=f'{rotation} scale(1.0, -1.0)')
        text_group.add(svg_text)
        return text_group


    def __get_colour(self, pptx_shape: PptxConnector | PptxGroupShape | PptxShape,
                           group_colour: Optional[ColourPair]=None) -> ColourPair:
    #=============================================================================
        def colour_from_fill(pptx_shape, fill) -> ColourPair:
            if fill.type == MSO_FILL_TYPE.SOLID:                                    # type: ignore
                return (self.__colour_map.lookup(fill.fore_color),
                        fill.fore_color.alpha)
            elif fill.type == MSO_FILL_TYPE.GRADIENT:                               # type: ignore
                log.warning(f'{pptx_shape.name}: gradient fill ignored')
            elif fill.type == MSO_FILL_TYPE.GROUP:                                  # type: ignore
                if group_colour is not None:
                    return group_colour
            elif fill.type is not None and fill.type != MSO_FILL_TYPE.BACKGROUND:   # type: ignore
                log.warning(f'{pptx_shape.name}: unsupported fill type: {fill.type}')
            return (None, 1.0)

        colour = None
        alpha = 1.0
        if pptx_shape.shape_type == MSO_SHAPE_TYPE.GROUP:                            # type: ignore
            colour, alpha = colour_from_fill(pptx_shape, FillFormat.from_fill_parent(pptx_shape.element.grpSpPr))
        elif pptx_shape.shape_type != MSO_SHAPE_TYPE.LINE:                           # type: ignore
            colour, alpha = colour_from_fill(pptx_shape, pptx_shape.fill)            # type: ignore
        elif pptx_shape.line.fill.type == MSO_FILL_TYPE.SOLID:                       # type: ignore
            colour = self.__colour_map.lookup(pptx_shape.line.color)                 # type: ignore
            alpha = pptx_shape.line.fill.fore_color.alpha                            # type: ignore
        elif pptx_shape.line.fill.type is None:                                      # type: ignore
            # Check for a fill colour in the <style> block
            xml = etree.fromstring(pptx_shape.element.xml)
            if (scheme_colour := xml.find('.//p:style/a:fillRef/a:schemeClr',
                                            namespaces=PPTX_NAMESPACE)) is not None:
                colour = self.__colour_map.scheme_colour(scheme_colour.attrib['val'])
        elif pptx_shape.line.fill.type != MSO_FILL_TYPE.BACKGROUND:                      # type: ignore
            log.warning(f'{pptx_shape.name}: unsupported line fill type: {pptx_shape.line.fill.type}')  # type: ignore
        return (colour, alpha)


    def __get_fill(self, pptx_shape: PptxConnector | PptxShape,
                         group_colour: Optional[ColourPair]=None) -> dict[str, Any]:
    #===============================================================================
        fill_attribs = {}
        colour, opacity = self.__get_colour(pptx_shape, group_colour)
        if (pptx_shape.fill.type == MSO_FILL_TYPE.SOLID                             # type: ignore
         or pptx_shape.fill.type == MSO_FILL_TYPE.GROUP):                           # type: ignore
            fill_attribs['fill'] = colour
            if opacity < 1.0:
                fill_attribs['opacity'] = opacity
        elif pptx_shape.fill.type == MSO_FILL_TYPE.GRADIENT:                        # type: ignore
            self.__gradient_id += 1
            gradient = Gradient(self.__drawing.defs, self.__gradient_id, pptx_shape, self.__colour_map)
            fill_attribs['fill'] = gradient.url
        elif pptx_shape.fill.type is None:                                          # type: ignore
            fill_attribs['fill'] = '#FF0000'
            fill_attribs['opacity'] = 1.0
        elif pptx_shape.fill.type != MSO_FILL_TYPE.BACKGROUND:                      # type: ignore
            print('Unsupported fill type: {}'.format(pptx_shape.fill.type))         # type: ignore
        return fill_attribs

    @staticmethod
    def __get_link(pptx_shape: PptxConnector | PptxShape) -> Optional[str]:
    #======================================================================
        shape_xml = etree.fromstring(pptx_shape.element.xml)
        for link_ref in shape_xml.findall('.//a:hlinkClick', namespaces=PPTX_NAMESPACE):
            r_id = link_ref.attrib[pptx_resolve('r:id')]
            if (r_id in pptx_shape.part.rels
             and pptx_shape.part.rels[r_id].reltype == pptx_uri('r:hyperlink')):
                return pptx_shape.part.rels[r_id].target_ref

    def __get_stroke(self, pptx_shape: PptxConnector | PptxShape) -> dict[str, Any]:
    #===============================================================================
        stroke_attribs = {}
        stroke_width = points_to_pixels(max(Length(pptx_shape.line.width).pt, MIN_STROKE_WIDTH))  # type: ignore
        stroke_attribs['stroke-width'] = stroke_width
        shape_xml = etree.fromstring(pptx_shape.element.xml)
        line_dash = pptx_shape.line.prstDash                                        # type: ignore
        try:
            dash_style = pptx_shape.line.dash_style                                 # type: ignore
        except KeyError:
            dash_style = None
        if line_dash is not None or dash_style is not None:
            if dash_style == MSO_LINE_DASH_STYLE.DASH:                              # type: ignore
                stroke_attribs['stroke-dasharray'] = 4*stroke_width
            elif line_dash == 'sysDot':
                stroke_attribs['stroke-dasharray'] = '{} {} {} {}'.format(4*stroke_width, stroke_width, stroke_width, stroke_width)
            elif line_dash == MSO_LINE_DASH_STYLE.LONG_DASH:                        # type: ignore
                stroke_attribs['stroke-dasharray'] = '{} {}'.format(4*stroke_width, stroke_width)
            elif dash_style == MSO_LINE_DASH_STYLE.SQUARE_DOT:                      # type: ignore
                stroke_attribs['stroke-dasharray'] = '{} {}'.format(2*stroke_width, stroke_width)
            elif dash_style == MSO_LINE_DASH_STYLE.ROUND_DOT:                       # type: ignore
                stroke_attribs['stroke-dasharray'] = '{} {}'.format(stroke_width, stroke_width)
            elif line_dash != 'solid':
                print(f'Unsupported line dash style: {dash_style}/{line_dash}')

        if pptx_shape.line.fill.type == MSO_FILL_TYPE.SOLID:                        # type: ignore
            stroke_attribs['stroke'] = self.__colour_map.lookup(pptx_shape.line.color)  # type: ignore
            alpha = pptx_shape.line.fill.fore_color.alpha                           # type: ignore
            if alpha < 1.0:
                stroke_attribs['stroke-opacity'] = alpha
        elif (line_style := shape_xml.find('.//p:style/a:lnRef', namespaces=PPTX_NAMESPACE)) is not None:
            for prop in line_style.getchildren():
                if prop.tag == DRAWINGML('schemeClr'):
                    scheme_colour = prop.attrib.get('val')
                    stroke_attribs['stroke'] = self.__colour_map.scheme_colour(scheme_colour)
        elif pptx_shape.line.fill.type is None:                                     # type: ignore
            stroke_attribs['stroke'] = 'none'
        elif pptx_shape.line.fill.type != MSO_FILL_TYPE.BACKGROUND:                 # type: ignore
            print('Unsupported line fill type: {}'.format(pptx_shape.line.fill.type))  # type: ignore
        return stroke_attribs

#===============================================================================

class SvgMaker:
    def __init__(self, powerpoint: Powerpoint, base_maker=None):
        if base_maker is None:
            self.__drawing = SvgDrawing(size=None)
            add_marker_definitions(self.__drawing)
            ## World --> pixels
            bounds = powerpoint.bounds   # southwest and northeast corners
            T = Transform(
                [[1.0, 0.0, -bounds[0]/WORLD_METRES_PER_PIXEL],
                 [0.0, 1.0 , bounds[3]/WORLD_METRES_PER_PIXEL],
                 [0.0, 0.0,                               1.0]])@np.array(
                    [[1/WORLD_METRES_PER_PIXEL,                       0.0, 0.0],
                     [                     0.0, -1/WORLD_METRES_PER_PIXEL, 0.0],
                     [                     0.0,                       0.0, 1.0]])
            svg_size = T.transform_point((bounds[2], bounds[1]))
            self.__drawing.attribs['viewBox'] = f'0 0 {svg_size[0]} {svg_size[1]}'
            self.__transform_matrix = svgelements.Matrix(T.svg_matrix)
            self.__celldl = CellDLGraph() if 'exportSVG' in settings else None
        else:
            self.__celldl = base_maker.__celldl
            self.__drawing = base_maker.__drawing
            self.__transform_matrix = base_maker.__transform_matrix

    def add_slides(self, slides: dict[tuple[int, str], Slide]):
    #==========================================================
        for (_, id), slide in slides.items():
            # in a group...
            slide_svg_maker = SvgFromSlide(id, slide, self.__drawing, self.__transform_matrix,
                                           celldl=self.__celldl)  ## flip_text
            slide_svg_maker.process_slide_svgs()

    def save(self, file_object):
    #===========================
        if self.__celldl is not None:
            self.__drawing.set_desc(desc='CellDL Metadata')
            self.__drawing.elements[0].xml.attrib['data-metadata-format'] = 'text/turtle'
            self.__drawing.elements[0].xml.attrib['data-metadata'] = self.__celldl.as_encoded_turtle()
        self.__drawing.write(file_object, pretty=True, indent=4)

    def svg_bytes(self):
    #===================
        svg = StringIO()
        self.__drawing.write(svg)
        svg_bytes = svg.getvalue().encode('utf-8')
        svg.close()
        return svg_bytes

#===============================================================================
