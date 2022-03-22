#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019 - 2022  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

import numpy as np
import pptx.shapes.connector
from pptx.enum.shapes import MSO_SHAPE_TYPE

#===============================================================================

from mapmaker.flatmap.layers import MapLayer
from mapmaker.utils import ProgressBar, log

from ..markup import parse_layer_directive

from .transform import DrawMLTransform
from .utils import get_shape_geometry

#===============================================================================

class PowerpointSlide(MapLayer):
    def __init__(self, source, slide, slide_number):
        id = 'slide-{:02d}'.format(slide_number)
        # Get any layer directives
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            if notes_text.startswith('.'):
                layer_directive = parse_layer_directive(notes_text)
                if 'error' in layer_directive:
                    source.error('error', 'Slide {}: invalid layer directive: {}'
                                 .format(slide_number, notes_text))
                if 'id' in layer_directive:
                    id = layer_directive['id']
        super().__init__(id, source, exported=(slide_number==1))
        self.__slide = slide
        self.__slide_number = slide_number
        self.__transform = source.transform

    @property
    def slide(self):
        return self.__slide

    @property
    def slide_id(self):
        return self.__slide.slide_id

    @property
    def slide_number(self):
        return self.__slide_number

    def process(self):
    #=================
        features = self.__process_shape_list(self.slide.shapes, self.__transform, show_progress=True)
        self.add_features('Slide', features, outermost=True)

    def __process_group(self, group, properties, transform):
    #=======================================================
        features = self.__process_shape_list(group.shapes, transform@DrawMLTransform(group))
        return self.add_features(properties.get('markup', ''), features)

    def __process_shape_list(self, shapes, transform, show_progress=False):
    #======================================================================
        progress_bar = ProgressBar(show=show_progress,
            total=len(shapes),
            unit='shp', ncols=40,
            bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')
        features = []
        for shape in shapes:
            properties = {'tile-layer': 'features'}   # Passed through to map viewer
            properties.update(self.source.properties_from_markup(shape.name))
            if 'error' in properties:
                pass
            elif 'invisible' in properties:
                pass
            elif 'path' in properties:
                pass
            elif (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or isinstance(shape, pptx.shapes.connector.Connector)):
                geometry = get_shape_geometry(shape, transform, properties)
                if geometry is not None:
                    feature = self.flatmap.new_feature(geometry, properties)
                    features.append(feature)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                grouped_feature = self.__process_group(shape, properties, transform)
                if grouped_feature is not None:
                    features.append(grouped_feature)
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                if 'id' in properties:
                    geometry = get_shape_geometry(shape, transform, properties)
                    if geometry is not None:
                        feature = self.flatmap.new_feature(geometry, properties)
                        features.append(feature)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pass
            else:
                log.warning('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))
            progress_bar.update(1)
        progress_bar.close()
        return features

#===============================================================================
