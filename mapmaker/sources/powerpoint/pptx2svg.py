#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019 - 2022 David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from collections import OrderedDict
from math import sqrt, sin, cos
import os
from pathlib import Path
import re
from typing import Any, Optional

#===============================================================================

from lxml import etree
import svgwrite
import svgwrite.gradients
from svgwrite.base import BaseElement as SvgElement
from svgwrite.container import Group as SvgGroup
from svgwrite.container import Hyperlink as SvgHyperlink
from svgwrite.text import Text as SvgText

#===============================================================================

from pptx import Presentation
from pptx.dml.fill import FillFormat
from pptx.enum.dml import MSO_FILL_TYPE, MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.util import Length

from pptx.shapes.autoshape import Shape as PptxShape
from pptx.shapes.connector import Connector as PptxConnector
from pptx.shapes.group import GroupShape as PptxGroupShape
from pptx.shapes.shapetree import GroupShapes as PptxGroupShapes
from pptx.shapes.shapetree import SlideShapes as PptxSlideShapes
from pptx.slide import Slide as PptxSlide

#===============================================================================

from mapmaker.geometry import Transform
from mapmaker.utils import FilePath, log, ProgressBar

from .colour import ColourMap, ColourTheme
from .geometry import get_shape_geometry
from .presets import DML
from .powerpoint import Shape, SHAPE_TYPE
from .transform import DrawMLTransform

#===============================================================================

PPTX_NAMESPACE = {
    'p': "http://schemas.openxmlformats.org/presentationml/2006/main",
    'a': "http://schemas.openxmlformats.org/drawingml/2006/main",
    'r': "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
}

def pptx_resolve(qname: str) -> str:
#===================================
    parts = qname.split(':', 1)
    if len(parts) == 2 and parts[0] in PPTX_NAMESPACE:
        return f'{{{PPTX_NAMESPACE[parts[0]]}}}{parts[1]}'
    return qname

#===============================================================================

# Internal PPT units are EMUs (English Metric Units)

EMU_PER_CM  = 360000
EMU_PER_IN  = 914400

POINTS_PER_IN = 72

# SVG pixel resolution
PIXELS_PER_IN = 96
EMU_PER_PIXEL = EMU_PER_IN/PIXELS_PER_IN

# Minimum width for a stroked path in points
MIN_STROKE_WIDTH = 0.5

TEXT_MARGINS = (4, 1)    # pixels

#===============================================================================

def text_alignment(shape):
#=========================
    para = shape.text_frame.paragraphs[0].alignment
    vertical = shape.text_frame.vertical_anchor
    return ('left' if para in [PP_ALIGN.LEFT, PP_ALIGN.DISTRIBUTE, PP_ALIGN.JUSTIFY, PP_ALIGN.JUSTIFY_LOW] else  # type: ignore
            'right' if para == PP_ALIGN.RIGHT else              # type: ignore
            'centre',
            'top' if vertical == MSO_ANCHOR.TOP else            # type: ignore
            'bottom' if vertical == MSO_ANCHOR.BOTTOM else      # type: ignore
            'middle')

def text_content(shape):
#=======================
    text = shape.text.replace('\n', ' ').replace('\xA0', ' ').replace('\v', ' ').strip() # Newline, non-breaking space, vertical-tab
    return text if text not in ['', '.'] else None

#===============================================================================

def emu_to_pixels(emu):
#======================
    return emu/EMU_PER_PIXEL

def points_to_pixels(pts):
#=========================
    return pts*PIXELS_PER_IN/POINTS_PER_IN

#===============================================================================

def ellipse_point(a, b, theta):
#==============================
    a_sin_theta = a*sin(theta)
    b_cos_theta = b*cos(theta)
    circle_radius = sqrt(a_sin_theta**2 + b_cos_theta**2)
    return (a*b_cos_theta/circle_radius, b*a_sin_theta/circle_radius)

#===============================================================================

ARROW_MARKERS = {
    'triangle-head': 'M 10 0 L 0 5 L 10 10 z',
    'triangle-tail': 'M 0 0 L 10 5 L 0 10 z'
}

## NB. Adobe Illustrator 2020 doesn't appear to support marker definitions in SVG

def add_marker_definitions(drawing):
#===================================
    # arrowhead markers (see https://developer.mozilla.org/en-US/docs/Web/SVG/Element/marker)
    # 18 Jan 2023: markers appear in Chrome with black fill; no markers in Firefox
    for id, path in ARROW_MARKERS.items():
        marker = drawing.marker(id=id,
                                viewBox="0 0 10 10",
                                refX="5", refY="5",
                                markerUnits="userSpaceOnUse",
                                markerWidth="6",
                                markerHeight="6",
                                orient="auto")
        marker.add(drawing.path(d=path))   ## , fill='context-stroke' is not supported by svgwrite
        drawing.defs.add(marker)

def marker_id(marker_def, end):
#==============================
    marker_type = marker_def.get('type')
    return ('#{}-{}'.format(marker_type, end)
            if marker_type is not None
            else None)

#===============================================================================

# Don't set a path id for default shape names

EXCLUDED_NAME_PREFIXES = [
    'Freeform',
    'Group',
    'Oval',
    'Star',
]

# Markup that has been deprecated

EXCLUDED_NAME_MARKUP = [
    '.siblings',
]

# Check to see if we have a valid name and encode it as an id

def valid_markup(name):
#======================
    if name not in EXCLUDED_NAME_MARKUP:
        for prefix in EXCLUDED_NAME_PREFIXES:
            if name.startswith(prefix):
                return False
        return True
    return False

def add_markup(element, markup):
#===============================
    if valid_markup(markup):
        element.set_desc(title=markup)

def add_class(xml, cls):
#=======================
    if 'class' not in xml.attribs:
        xml.attribs['class'] = cls
    else:
        xml.attribs['class'] += f' {cls}'

#===============================================================================

class Gradient(object):
    def __init__(self, dwg, id, shape, colour_map):
        self.__id = 'gradient-{}'.format(id)
        fill = shape.fill
        gradient = None
        if fill._fill._gradFill.path is None:
            gradient = svgwrite.gradients.LinearGradient(id=self.__id)
            rotation = fill.gradient_angle
            if ('rotWithShape' in fill._fill._gradFill.attrib
             and fill._fill._gradFill.attrib['rotWithShape'] == '1'):
                rotation += shape.rotation
            if rotation != 0:
                gradient.rotate(rotation % 360, (0.5, 0.5))

        elif fill._fill._gradFill.path.attrib['path'] == 'circle':
                fill_to = fill._fill._gradFill.path.find('{http://schemas.openxmlformats.org/drawingml/2006/main}fillToRect').attrib
                tileRect = fill._fill._gradFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}tileRect')
                tile = tileRect.attrib if tileRect is not None else {}
                cx = (float(fill_to['l']) if 'l' in fill_to else float(fill_to['r']) + float(tile.get('l', 0.0)))/100000.0
                cy = (float(fill_to['t']) if 't' in fill_to else float(fill_to['b']) + float(tile.get('t', 0.0)))/100000.0
                sx = (float(fill_to['r']) if 'r' in fill_to else float(fill_to['l']) + float(tile.get('r', 0.0)))/100000.0
                sy = (float(fill_to['b']) if 'b' in fill_to else float(fill_to['t']) + float(tile.get('b', 0.0)))/100000.0
                if shape.width > shape.height:
                    scale_x = shape.height/shape.width
                    scale_y = 1.0
                elif shape.width < shape.height:
                    scale_x = 1.0
                    scale_y = shape.width/shape.height
                else:
                    scale_x = 1.0
                    scale_y = 1.0
                if len(tile) == 0:
                    radius = 1.0
                    if (cx, cy) != (0.5, 0.5) and (sx, sy) != (0.5, 0.5):
                        print('Preset radial gradient for shape:', shape.name)
                else:
                    radius = sqrt(((cx-sx)/scale_x)**2 + ((cy-sy)/scale_y)**2)
                gradient = svgwrite.gradients.RadialGradient((cx/scale_x, cy/scale_y), radius, id=self.__id)
                if shape.rotation != 0:
                    gradient.rotate(shape.rotation, (0.5, 0.5))
                if shape.width != shape.height:
                    gradient.scale(scale_x, scale_y)

        elif fill._fill._gradFill.path.attrib['path'] == 'rect':
            print('Rect fill ignored for', shape.name)
            return

        if gradient is not None:
            for stop in sorted(fill.gradient_stops, key=lambda stop: stop.position):
                gradient.add_stop_color(offset=stop.position,
                    color=colour_map.lookup(stop.color),
                    opacity=stop.color.alpha)
            dwg.defs.add(gradient)
        else:
            print('UNKNOWN FILL: {}\n'.format(shape.name), fill._fill._element.xml)

    @property
    def url(self):
        return 'url(#{})'.format(self.__id)

## WIP  Want list of unique gradient definitions

#===============================================================================

# (colour, opacity)
ColourPair = tuple[Optional[str], float]

#===============================================================================

class SvgLayer(object):
    def __init__(self, size, slide: PptxSlide, slide_number: int, ppt_theme, kind: str='base', shape_filter=None, quiet=False):
        self.__slide = slide
        self.__colour_map = ColourMap(ppt_theme, slide)
        self.__dwg = svgwrite.Drawing(size=None)
        self.__dwg.attribs['viewBox'] = f'0 0 {size[0]} {size[1]}'
        add_marker_definitions(self.__dwg)
        self.__id = None
        self.__models = None
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            if notes_text.startswith('.'):
                for part in notes_text[1:].split():
                    id_match = re.match(r'id *\((.*)\)', part)
                    if id_match is not None:
                        self.__id = id_match[1].strip()
                    models_match = re.match(r'models *\((.*)\)', part)
                    if models_match is not None:
                        self.__models = models_match[1].strip()
        if self.__id is None:
            self.__id = 'slide-{:02d}'.format(slide_number)
        self.__gradient_id = 0
        self.__kind = kind
        self.__shape_filter = shape_filter
        self.__quiet =  quiet

    @property
    def id(self):
        return self.__id

    @property
    def models(self):
        return self.__models

    def save(self, file_object):
    #===========================
        self.__dwg.write(file_object, pretty=True, indent=4)

    def __get_colour(self, shape: PptxConnector | PptxGroupShape | PptxShape,
                     group_colour: Optional[ColourPair]=None) -> ColourPair:
    #=======================================================================
        def colour_from_fill(shape, fill) -> ColourPair:
            if fill.type == MSO_FILL_TYPE.SOLID:                                    # type: ignore
                return (self.__colour_map.lookup(fill.fore_color),
                        fill.fore_color.alpha)
            elif fill.type == MSO_FILL_TYPE.GRADIENT:                               # type: ignore
                log.warning(f'{shape.name}: gradient fill ignored')
            elif fill.type == MSO_FILL_TYPE.GROUP:                                  # type: ignore
                if group_colour is not None:
                    return group_colour
            elif fill.type is not None and fill.type != MSO_FILL_TYPE.BACKGROUND:   # type: ignore
                log.warning(f'{shape.name}: unsupported fill type: {fill.type}')
            return (None, 1.0)

        colour = None
        alpha = 1.0
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:                                # type: ignore
            colour, alpha = colour_from_fill(shape, FillFormat.from_fill_parent(shape.element.grpSpPr))
        elif shape.shape_type != MSO_SHAPE_TYPE.LINE:                               # type: ignore
            colour, alpha = colour_from_fill(shape, shape.fill)                     # type: ignore
        elif shape.line.fill.type == MSO_FILL_TYPE.SOLID:                           # type: ignore
            colour = self.__colour_map.lookup(shape.line.color)                     # type: ignore
            alpha = shape.line.fill.fore_color.alpha                                # type: ignore
        elif shape.line.fill.type is None:                                          # type: ignore
            # Check for a fill colour in the <style> block
            xml = etree.fromstring(shape.element.xml)
            if (scheme_colour := xml.find('.//p:style/a:fillRef/a:schemeClr',
                                            namespaces=PPTX_NAMESPACE)) is not None:
                colour = self.__colour_map.scheme_colour(scheme_colour.attrib['val'])
        elif shape.line.fill.type != MSO_FILL_TYPE.BACKGROUND:                      # type: ignore
            log.warning(f'{shape.name}: unsupported line fill type: {shape.line.fill.type}')  # type: ignore
        return (colour, alpha)

    def process(self, transform: Transform):
    #=======================================
        self.process_shape_list(self.__slide.shapes,                                # type: ignore
                                self.__dwg, transform, show_progress=not self.__quiet)
        if self.__kind == 'base' and self.__shape_filter is not None:
            self.__shape_filter.create_filter()

    def process_group(self, group: PptxGroupShape, svg_parent: SvgElement, transform: Transform):
    #============================================================================================
        svg_group = SvgGroup(id=group.shape_id)
        add_markup(svg_group, group.name)
        svg_parent.add(svg_group)
        self.process_shape_list(group.shapes, svg_group,                            # type: ignore
                                transform@DrawMLTransform(group), group_colour=self.__get_colour(group))

    def process_shape_list(self, shapes: PptxGroupShapes | PptxSlideShapes, svg_parent: SvgElement,
                           transform: Transform, group_colour: Optional[ColourPair]=None,
                           show_progress: bool=False):
    #===============================================================================================
        if show_progress:
            print('Processing shape list...')
        progress_bar = ProgressBar(show=show_progress,
            total=len(shapes),
            unit='shp', ncols=40,
            bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')
        for shape in shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE                   # type: ignore
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM                     # type: ignore
             or shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX                     # type: ignore
             or shape.shape_type == MSO_SHAPE_TYPE.LINE):                       # type: ignore
                self.process_shape(shape, svg_parent, transform,                # type: ignore
                                   group_colour=group_colour)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:                      # type: ignore
                self.process_group(shape, svg_parent, transform)                # type: ignore
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:                    # type: ignore
                pass
            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))
            progress_bar.update(1)
        progress_bar.close()

    def process_shape(self, shape: PptxConnector | PptxShape, svg_parent: SvgElement,
                      transform: Transform, group_colour: Optional[ColourPair]=None):
    #================================================================================

        properties = {}
        geometry = get_shape_geometry(shape, transform, properties)
        svg_path = properties['svg-path']

        bbox = (shape.width, shape.height)
        T = transform@DrawMLTransform(shape, bbox)
        shape_size = T.scale_length(bbox)

        exclude_shape = False
        exclude_text = True
        svg_text = None
        label = None
        colour, opacity = self.__get_colour(shape, group_colour)
        if shape.shape_type != MSO_SHAPE_TYPE.LINE:                                 # type: ignore
            svg_path.attribs.update(self.__get_fill(shape, group_colour))
            label = text_content(shape)
            if self.__shape_filter is not None and properties.get('closed'):
                # Apply filter to closed shapes
                ## this doesn't work for shapes drawn using arcTo or *BezTo commands
                ## OK for FC since (most) shapes are rectangular (but not circular nodes...)
                shape_properties = {
                    'colour': colour,
                    'opacity': opacity
                }
                if label is not None:
                    shape_properties['label'] = label
                closed_shape = Shape(SHAPE_TYPE.FEATURE, shape.shape_id, geometry, shape_properties)
                if self.__kind == 'base':
                    self.__shape_filter.add_shape(closed_shape)
                elif self.__kind == 'layer':
                    self.__shape_filter.filter(closed_shape)
                exclude_shape = shape_properties.get('exclude', False)
                exclude_text = shape_properties.get('exclude-text', False)
            else:
                exclude_text = False

            exclude_text = True ####################### TEMP
            if (not exclude_text and not exclude_shape   ## <<<<<<<<<<<<<<<<<<<<< No text at all...
            and self.__kind == 'base' and label is not None):
                svg_text = self.__draw_shape_label(shape, label, shape_size, T)

        elif True:  #####  WIP self.__shape_filter is not None:   ## No filter for pptx2celldel ??
            # Exclude connectors when filtering shapes
            exclude_shape = True

            if 'type' in shape.line.headEnd or 'type' in shape.line.tailEnd:
                svg_path.set_markers((marker_id(shape.line.headEnd, 'head'),
                                      None, marker_id(shape.line.tailEnd, 'tail')
                                    ))

        if not exclude_shape:
            svg_path.attribs.update(self.__get_stroke(shape))

            shape_kind = properties['shape-kind']
            if (shape_kind is not None
             and not shape_kind.startswith('star')
             and shape_size[0]*shape_size[1] < 200):
                add_class(svg_path, 'connector')    ## FC
                                                    ## cardio (circle) versus neural (rect)
                                                    ## nerve features...

            if (hyperlink := self.__get_link(shape)) is not None:
                if label is None:
                    label = hyperlink
                link_element = SvgHyperlink(href=hyperlink)
                link_element.add(svg_path)
                if svg_text is not None:
                    link_element.add(svg_text)
                svg_parent.add(link_element)
            else:
                svg_parent.add(svg_path)
                if svg_text is not None:
                    svg_parent.add(svg_text)
            if label is not None:
                add_markup(svg_path, label)  # Set's <title>
                if svg_text is not None:
                    add_markup(svg_text, label)
            else:
                add_markup(svg_path, shape.name)

    def __draw_shape_label(self, shape: PptxShape, label: str, shape_size: tuple[float, float], transform: Transform) -> SvgElement:
    #===============================================================================================================================
        # Draw text if base map

        ##for paragraph in shape.text_frame.paragraphs:
        ##  for text_run in paragraph.runs:

        style = {}
        font = shape.text_frame.paragraphs[0].runs[0].font          # type: ignore
        font_size = round(font.size/EMU_PER_PIXEL)                  # type: ignore
        if font.name is not None:   ## Else need to get from theme
            style['font-family'] = font.name
        style['font-size'] = f'{font_size}px'
        style['font-weight'] = 700 if font.bold else 400
        if font.italic:
            style['font-size'] = 'italic'
        if font.color.type is not None:                             # type: ignore
            style['fill'] = self.__colour_map.lookup(font.color)
            if font.color.alpha != 1.0:                             # type: ignore
                style['fill-opacity'] = font.color.alpha            # type: ignore

        svg_text = SvgText(label)   ## text_run.text
        svg_text.attribs['style'] = ' '.join([f'{name}: {value};' for name, value in style.items()])

        shape_pos = transform.transform_point((0, 0))
        svg_text.attribs['x'] = shape_pos[0]
        svg_text.attribs['y'] = shape_pos[1]
        (halign, valign) = text_alignment(shape)

        if halign == 'right':
            svg_text.attribs['text-anchor'] = 'end'
            svg_text.attribs['x'] += shape_size[0] - TEXT_MARGINS[0]
        elif halign == 'centre':
            svg_text.attribs['text-anchor'] = 'middle'
            svg_text.attribs['x'] += shape_size[0]/2
        else:   # Default to 'left'
            svg_text.attribs['text-anchor'] = 'start'
            svg_text.attribs['x'] += TEXT_MARGINS[0]
        if valign == 'bottom':
            svg_text.attribs['dominant-baseline'] = 'auto'
            svg_text.attribs['y'] += shape_size[1]
        elif valign == 'middle':
            svg_text.attribs['dominant-baseline'] = 'middle'
            svg_text.attribs['y'] += shape_size[1]/2
        else:   # Default to 'top'
            svg_text.attribs['dominant-baseline'] = 'auto'
            svg_text.attribs['y'] += font_size + TEXT_MARGINS[1]
        return svg_text

    def __get_fill(self, shape: PptxConnector | PptxShape, group_colour: Optional[ColourPair]=None) -> dict[str, Any]:
    #=================================================================================================================
        fill_attribs = {}
        colour, opacity = self.__get_colour(shape, group_colour)
        if (shape.fill.type == MSO_FILL_TYPE.SOLID                      # type: ignore
         or shape.fill.type == MSO_FILL_TYPE.GROUP):                    # type: ignore
            fill_attribs['fill'] = colour
            if opacity < 1.0:
                fill_attribs['opacity'] = opacity
        elif shape.fill.type == MSO_FILL_TYPE.GRADIENT:                 # type: ignore
            self.__gradient_id += 1
            gradient = Gradient(self.__dwg, self.__gradient_id, shape, self.__colour_map)
            fill_attribs['fill'] = gradient.url
        elif shape.fill.type is None:                                   # type: ignore
            fill_attribs['fill'] = '#FF0000'
            fill_attribs['opacity'] = 1.0
        elif shape.fill.type != MSO_FILL_TYPE.BACKGROUND:               # type: ignore
            print('Unsupported fill type: {}'.format(shape.fill.type))  # type: ignore
        return fill_attribs

    @staticmethod
    def __get_link(shape: PptxConnector | PptxShape) -> Optional[str]:
    #=================================================================
        shape_xml = etree.fromstring(shape.element.xml)
        for link_ref in shape_xml.findall('.//a:hlinkClick', namespaces=PPTX_NAMESPACE):
            r_id = link_ref.attrib[pptx_resolve('r:id')]
            if (r_id in shape.part.rels
             and shape.part.rels[r_id].reltype == 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink'):
                return shape.part.rels[r_id].target_ref

    def __get_stroke(self, shape: PptxConnector | PptxShape) -> dict[str, Any]:
    #==========================================================================
        stroke_attribs = {}
        stroke_width = points_to_pixels(max(Length(shape.line.width).pt, MIN_STROKE_WIDTH))  # type: ignore
        stroke_attribs['stroke-width'] = stroke_width
        shape_xml = etree.fromstring(shape.element.xml)

        line_dash = None
        if (line_props := shape_xml.find('.//p:spPr/a:ln', namespaces=PPTX_NAMESPACE)) is not None:
            for prop in line_props.getchildren():
                if prop.tag == DML('prstDash'):
                    line_dash = prop.attrib.get('val', 'solid')
                    break
        try:
            dash_style = shape.line.dash_style                          # type: ignore
        except KeyError:
            dash_style = None
        if line_dash is not None or dash_style is not None:
            if dash_style == MSO_LINE_DASH_STYLE.DASH:                  # type: ignore
                stroke_attribs['stroke-dasharray'] = 4*stroke_width
            elif line_dash == 'sysDot':
                stroke_attribs['stroke-dasharray'] = '{} {} {} {}'.format(4*stroke_width, stroke_width, stroke_width, stroke_width)
            elif line_dash == MSO_LINE_DASH_STYLE.LONG_DASH:            # type: ignore
                stroke_attribs['stroke-dasharray'] = '{} {}'.format(4*stroke_width, stroke_width)
            elif dash_style == MSO_LINE_DASH_STYLE.SQUARE_DOT:          # type: ignore
                stroke_attribs['stroke-dasharray'] = '{} {}'.format(2*stroke_width, stroke_width)
            elif dash_style == MSO_LINE_DASH_STYLE.ROUND_DOT:           # type: ignore
                stroke_attribs['stroke-dasharray'] = '{} {}'.format(stroke_width, stroke_width)
            elif line_dash != 'solid':
                print(f'Unsupported line dash style: {dash_style}/{line_dash}')

        if shape.line.fill.type == MSO_FILL_TYPE.SOLID:                 # type: ignore
            stroke_attribs['stroke'] = self.__colour_map.lookup(shape.line.color)  # type: ignore
            alpha = shape.line.fill.fore_color.alpha                    # type: ignore
            if alpha < 1.0:
                stroke_attribs['stroke-opacity'] = alpha
        elif (line_style := shape_xml.find('.//p:style/a:lnRef', namespaces=PPTX_NAMESPACE)) is not None:
            for prop in line_style.getchildren():
                if prop.tag == DML('schemeClr'):
                    scheme_colour = prop.attrib.get('val')
                    stroke_attribs['stroke'] = self.__colour_map.scheme_colour(scheme_colour)
        elif shape.line.fill.type is None:                              # type: ignore
            stroke_attribs['stroke'] = 'none'
        elif shape.line.fill.type != MSO_FILL_TYPE.BACKGROUND:          # type: ignore
            print('Unsupported line fill type: {}'.format(shape.line.fill.type))  # type: ignore

        return stroke_attribs

#===============================================================================

class Pptx2Svg(object):
    def __init__(self, powerpoint_href, kind='base', shape_filter=None, quiet=True):
        self.__source_name = Path(powerpoint_href).stem
        ppt_bytes = FilePath(powerpoint_href).get_BytesIO()
        self.__pptx = Presentation(ppt_bytes)
        self.__colour_theme = ColourTheme(ppt_bytes)
        self.__slides = self.__pptx.slides
        (pptx_width, pptx_height) = (self.__pptx.slide_width, self.__pptx.slide_height)
        self.__transform = Transform([[1.0/EMU_PER_PIXEL,                 0, 0],
                                      [                0, 1.0/EMU_PER_PIXEL, 0],
                                      [                0,                 0, 1]])
        self.__svg_size = self.__transform.transform_point((pptx_width, pptx_height))
        self.__kind = kind
        self.__shape_filter = shape_filter
        self.__quiet = quiet
        self.__svg_layers = []
        self.__saved_svg = OrderedDict()
        self.__id = Path(powerpoint_href).name.split('.')[0].replace(' ', '_')
        self.__models = None

    @property
    def id(self):
        return self.__id

    @property
    def svg_layers(self):
        return self.__svg_layers

    def slide_to_svg(self, slide, slide_number):
    #===========================================
        layer = SvgLayer(self.__svg_size, slide, slide_number, self.__colour_theme,
                         kind=self.__kind, shape_filter=self.__shape_filter,
                         quiet=self.__quiet)
        layer.process(self.__transform)
        self.__svg_layers.append(layer)
        if slide_number == 1:
            if layer.id is not None and not layer.id.startswith('slide-'):
                self.__id = layer.id
            self.__models = layer.models

    def slides_to_svg(self):
    #=======================
        for n, slide in enumerate(self.__slides):
            self.slide_to_svg(slide, n+1)

    def save_layers(self, output_dir):
    #=================================
        self.__saved_svg = OrderedDict()
        layer_ids = len(self.__svg_layers) > 1
        for layer in self.__svg_layers:
            if layer_ids:
                filename = f'{self.__source_name}-{layer.id}.svg'
            else:
                filename = f'{self.__source_name}.svg'
            svg_file = os.path.join(output_dir, filename)
            with open(svg_file, 'w', encoding='utf-8') as fp:
                layer.save(fp)
            self.__saved_svg[layer.id] = svg_file
        return self.__saved_svg

    def update_manifest(self, manifest):
    #===================================
        if 'id' not in manifest and self.__id is not None:
            manifest['id'] = self.__id
        if 'models' not in manifest and self.__models is not None:
            manifest['models'] = self.__models
        sources = [ source for source in manifest['sources']
                        if source['kind'] != 'slides']
        next_kind = 'base'
        for id, filename in self.__saved_svg.items():
            sources.append(OrderedDict(
                id=id,
                href=filename,
                kind=next_kind
            ))
            if next_kind == 'base':
                next_kind = 'details'
        manifest['sources'] = sources

#===============================================================================
