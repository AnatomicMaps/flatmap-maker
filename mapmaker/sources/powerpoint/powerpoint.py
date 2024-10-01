#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2018 - 2023  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

import base64
from typing import Optional, TYPE_CHECKING

#===============================================================================

import lxml.etree as etree
import numpy as np
import shapely.geometry
from shapely.geometry.base import BaseGeometry
import shapely.ops
import svgelements

from pptx import Presentation
from pptx.dml.fill import FillFormat
from pptx.enum.dml import MSO_FILL_TYPE             # type: ignore
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN     # type: ignore
from pptx.oxml.text import CT_RegularTextRun, CT_TextLineBreak, CT_TextField
from pptx.shapes.autoshape import Shape as PptxShape
from pptx.shapes.connector import Connector as PptxConnector
from pptx.shapes.group import GroupShape as PptxGroupShape
from pptx.shapes.shapetree import GroupShapes as PptxGroupShapes
from pptx.shapes.shapetree import SlideShapes as PptxSlideShapes
from pptx.slide import Slide as PptxSlide

#===============================================================================

from mapmaker.geometry import Transform
from mapmaker.properties.markup import parse_layer_directive, parse_markup
from mapmaker.shapes import Shape, SHAPE_TYPE
from mapmaker.sources import MapBounds, WORLD_METRES_PER_EMU
from mapmaker.utils import FilePath, log, ProgressBar, TreeList

from ..fc_powerpoint.colours import ColourMatcher
from ..fc_powerpoint.components import is_system_name

from .colour import ColourMap, ColourTheme
from .geometry import get_shape_geometry
from .presets import CT_TextMath, DRAWINGML, PPTX_NAMESPACE, pptx_resolve, pptx_uri
from .transform import DrawMLTransform

from .omml2latex import openmath2latex

if TYPE_CHECKING:
    from mapmaker.annotation import Annotator

#===============================================================================

STROKE_WIDTH_SCALE_FACTOR = 1270.0

#===============================================================================

# (colour, opacity)
ColourPair = tuple[Optional[str], float]

#===============================================================================

class Slide:
    def __init__(self, flatmap: 'FlatMap', source: 'PowerpointSource', index: int, pptx_slide: PptxSlide,   # type: ignore
                 theme: ColourTheme, bounds: MapBounds, transform: Transform):
        self.__flatmap = flatmap
        self.__source = source
        self.__kind = source.kind
        self.__id = 'slide-{:02d}'.format(index+1)
        # Get any layer directives
        if pptx_slide.has_notes_slide:
            notes_slide = pptx_slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            if notes_text.startswith('.'):
                layer_directive = parse_layer_directive(notes_text)
                if 'error' in layer_directive:
                    log.error('error', f'Slide {index+1}: invalid layer directive: {notes_text}')
                if 'id' in layer_directive:
                    self.__id = layer_directive['id']
        self.__colour_map = ColourMap(theme, pptx_slide)
        self.__pptx_slide = pptx_slide
        self.__geometry = shapely.geometry.box(*bounds)
        self.__transform = transform
        self.__shapes = TreeList()
        self.__shapes_by_id: dict[str, Shape] = {}

    @property
    def colour_map(self) -> ColourMap:
        return self.__colour_map

    @property
    def flatmap(self):
        return self.__flatmap

    @property
    def geometry(self) -> shapely.geometry.base.BaseGeometry:
        return self.__geometry

    @property
    def kind(self) -> str:
        return self.__kind

    @property
    def id(self) -> str:
        return self.__id

    @property
    def pptx_slide(self) -> PptxSlide:
        return self.__pptx_slide

    @property
    def shapes(self) -> TreeList:
        return self.__shapes

    @property
    def slide_id(self) -> int:
        return self.__pptx_slide.slide_id

    @property
    def source(self) -> 'PowerpointSource':     # type: ignore
        return self.__source

    @property
    def source_id(self) -> str:
        return self.__source.id

    def shape(self, id: str) -> Optional[Shape]:
    #===========================================
        return self.__shapes_by_id.get(id)

    def __shape_id(self, id) -> str:
    #==============================
        return f'{self.__source.id}/{self.__id}/{id}'

    def __new_shape(self, id: str, geometry, properties, shape_type=None) -> Shape:
    #==============================================================================
        shape_id = self.__shape_id(id)
        shape = Shape(shape_id, geometry, properties, shape_type=shape_type)
        self.__shapes_by_id[shape_id] = shape
        return shape

    def process(self, annotator: Optional['Annotator']=None) -> TreeList:
    #==================================================================
        # Return the slide's group structure as a nested list of Shapes
        self.__shapes = TreeList([self.__new_shape('root', self.__geometry, {'type': SHAPE_TYPE.GROUP})])
        self.__shapes.extend(self.__process_pptx_shapes(self.__pptx_slide.shapes,      # type: ignore
                                                        self.__transform, show_progress=True))
        return self.__shapes

    def __get_colour(self, shape: PptxConnector | PptxGroupShape | PptxShape,
                     group_colour: Optional[ColourPair]=None) -> ColourPair:
    #=======================================================================
        def colour_from_fill(shape, fill) -> ColourPair:
            if fill.type == MSO_FILL_TYPE.SOLID:                    # type: ignore
                return (self.__colour_map.lookup(fill.fore_color),
                        fill.fore_color.alpha)
            elif fill.type == MSO_FILL_TYPE.GRADIENT:               # type: ignore
                colours = [(self.__colour_map.lookup(stop.color), stop.color.alpha)
                                for stop in fill.gradient_stops]
                n = 0
                while n < len(colours) and colours[n][0] == '#FFFFFF':
                    n += 1
                if n >= len(colours):
                    n = 0
                log.warning(f'{shape.text}: gradient fill ignored, stop colour `{n}` used')
                return colours[n]
            elif fill.type == MSO_FILL_TYPE.GROUP:                  # type: ignore
                if group_colour is not None:
                    return group_colour
            elif fill.type is not None and fill.type != MSO_FILL_TYPE.BACKGROUND:   # type: ignore
                log.warning(f'{shape.text}: unsupported fill type: {fill.type}')
            return (None, 1.0)

        colour = None
        alpha = 1.0
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:                # type: ignore
            colour, alpha = colour_from_fill(shape, FillFormat.from_fill_parent(shape.element.grpSpPr))
        elif shape.shape_type != MSO_SHAPE_TYPE.LINE:               # type: ignore
            colour, alpha = colour_from_fill(shape, shape.fill)     # type: ignore
        elif shape.line.fill.type == MSO_FILL_TYPE.SOLID:           # type: ignore
            colour = self.__colour_map.lookup(shape.line.color)     # type: ignore
            alpha = shape.line.fill.fore_color.alpha                # type: ignore
        elif shape.line.fill.type is None:                          # type: ignore
            # Check for a fill colour in the <style> block
            xml = etree.fromstring(shape.element.xml)
            if (scheme_colour := xml.find('.//p:style/a:fillRef/a:schemeClr',
                                            namespaces=PPTX_NAMESPACE)) is not None:
                colour = self.__colour_map.scheme_colour(scheme_colour.attrib['val'])
        elif shape.line.fill.type != MSO_FILL_TYPE.BACKGROUND:      # type: ignore
            log.warning(f'{shape.text}: unsupported line fill type: {shape.line.fill.type}')    # type: ignore
        return (colour, alpha)

    def __shapes_as_group(self, group: PptxGroupShape, shapes: TreeList) -> Shape | TreeList:
    #========================================================================================
        # Merge a group of overlapping shapes that are all the same colour
        # and have a common name into a single shape
        if (len(shapes) < 2
         or isinstance(shapes[0], TreeList)
         or shapes[0].shape_type != SHAPE_TYPE.COMPONENT):
            return shapes
        name = shapes[0].name
        colour = ColourMatcher(shapes[0].colour)
        alignment = shapes[0].properties.get('align')
        geometries = [shapes[0].geometry]
        pptx_shape = shapes[0].properties['pptx-shape']
        for shape in shapes[1:]:
            if (isinstance(shape, TreeList)
             or shape.shape_type != SHAPE_TYPE.COMPONENT
             or not colour.matches(shape.colour)):
                return shapes
            if shape.name != '':
                if name == '':
                    name = shape.name
                    colour = ColourMatcher(shape.colour)
                    alignment = shape.properties.get('align')
                    pptx_shape = shape.properties['pptx-shape']
                elif name != shape.name:
                    return shapes
            geometries.append(shape.geometry)

        if name == '':
            return shapes
        geometry = shapely.ops.unary_union(geometries)
        if ('Polygon' not in geometry.geom_type
         or geometry.geom_type != 'Polygon' and not is_system_name(name)):
            return shapes

        svg = etree.fromstring(geometry.svg())
        svg_elements = []
        if svg.tag == 'path' and (svg_path := svg.attrib.get('d')) is not None:
            svg_elements.append(svgelements.Path(svg_path))
        elif svg.tag == 'g':
            for e in svg:
                if e.tag == 'path' and (svg_path := e.attrib.get('d')) is not None:
                    svg_elements.append(svgelements.Path(svg_path))
        if len(svg_elements) == 0:
            return shapes

        return self.__new_shape(group.shape_id, geometry, {
                                'colour': colour.rgb_colour,
                                'name': name,
                                'shape-name': group.name,
                                'shape-kind': shapes[0].kind,
                                'text-align': alignment,
                                'pptx-shape': pptx_shape,
                                'svg-element': svg_elements[0] if len(svg_elements) == 1 else svg_elements,
                                'svg-kind': 'path' if len(svg_elements) == 1 else 'group'
                                }, SHAPE_TYPE.COMPONENT)

    def __text_content(self, shape: PptxShape) -> str:
    #=================================================
        text_types = {CT_RegularTextRun, CT_TextLineBreak, CT_TextField}
        shape_text = []
        for paragraph in shape.text_frame.paragraphs:
            paragraph_text = []
            for child in paragraph._element:
                if type(child) in text_types:
                    paragraph_text.append(child.text.replace('\n', ' ')
                                                    .replace('\xA0', ' ')
                                                    .replace('\v', ' '))        # Newline, non-breaking space, vertical-tab
                elif isinstance(child, CT_TextMath):
                    xml = etree.tostring(child.getchildren()[0], encoding='unicode')
                    latex = openmath2latex(xml)
                    paragraph_text.append(f'`{latex}`')
            shape_text.append(''.join(paragraph_text))
        text = ' '.join(shape_text)
        return ' '.join(text.split()) if text not in ['', '.'] else ''

    def __process_group(self, group: PptxGroupShape, transform: Transform) -> Shape | TreeList:
    #==========================================================================================
        colour = self.__get_colour(group)
        group_shapes = self.__shapes_as_group(group,
                            self.__process_pptx_shapes(group.shapes,        # type: ignore
                                transform@DrawMLTransform(group),
                                group_colour=colour))
        if isinstance(group_shapes, Shape):
            return group_shapes
        shapes = TreeList([self.__new_shape(group.shape_id, None, {
            'colour': colour[0],
            'opacity': colour[1],
            'pptx-shape': group
        }, SHAPE_TYPE.GROUP)])
        shapes.extend(group_shapes)
        return shapes

    def __process_pptx_shapes(self, pptx_shapes: PptxGroupShapes | PptxSlideShapes,
                              transform: Transform, group_colour: Optional[ColourPair]=None,
                              show_progress=False) -> TreeList:
    #========================================================================================
        def text_alignment(shape) -> tuple[str, str]:
            para = shape.text_frame.paragraphs[0].alignment
            vertical = shape.text_frame.vertical_anchor
            return ('left' if para in [PP_ALIGN.LEFT, PP_ALIGN.DISTRIBUTE, PP_ALIGN.JUSTIFY, PP_ALIGN.JUSTIFY_LOW] else
                    'right' if para == PP_ALIGN.RIGHT else
                    'centre',
                    'top' if vertical == MSO_ANCHOR.TOP else
                    'bottom' if vertical == MSO_ANCHOR.BOTTOM else
                    'middle')

        progress_bar = ProgressBar(show=show_progress,
            total=len(pptx_shapes),
            unit='shp', ncols=40,
            bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')
        shapes = TreeList()
        for pptx_shape in pptx_shapes:
            shape_name = pptx_shape.name
            shape_properties = parse_markup(shape_name) if shape_name.startswith('.') else {}
            shape_properties['pptx-shape'] = pptx_shape
            shape_properties['shape-name'] = shape_name

            def good_geometry(geometry):
                if geometry is None:
                    log.warning(f'Shape "{shape_name}" {pptx_shape.shape_type}/{shape_properties.get("shape-kind")} not processed -- cannot get geometry')
                elif not geometry.is_valid:
                    log.warning(f'Shape "{shape_name}" {pptx_shape.shape_type}/{shape_properties.get("shape-kind")} not processed -- cannot get valid geometry')
                else:
                    return True
                return False

            if (pptx_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE              # type: ignore
             or pptx_shape.shape_type == MSO_SHAPE_TYPE.FREEFORM                # type: ignore
             or pptx_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX                # type: ignore
             or pptx_shape.shape_type == MSO_SHAPE_TYPE.LINE):                  # type: ignore
                colour, alpha = self.__get_colour(pptx_shape, group_colour)     # type: ignore
                shape_properties['colour'] = colour
                if alpha < 1.0:
                    shape_properties['opacity'] = alpha
                if good_geometry(geometry := get_shape_geometry(pptx_shape, transform, shape_properties)):
                    shape_xml = etree.fromstring(pptx_shape.element.xml)
                    for link_ref in shape_xml.findall('.//a:hlinkClick',
                                                    namespaces=PPTX_NAMESPACE):
                        r_id = link_ref.attrib[pptx_resolve('r:id')]
                        if (r_id in pptx_shape.part.rels
                         and pptx_shape.part.rels[r_id].reltype == pptx_uri('r:hyperlink')):
                            shape_properties['hyperlink'] = pptx_shape.part.rels[r_id].target_ref
                            break
                    if pptx_shape.shape_type == MSO_SHAPE_TYPE.LINE:            # type: ignore
                        ## cf. pptx2svg for stroke colour
                        shape_type = SHAPE_TYPE.CONNECTION
                        if (connection := shape_xml.find('.//p:nvCxnSpPr/p:cNvCxnSpPr',
                                                        namespaces=PPTX_NAMESPACE)) is not None:
                            for c in connection.getchildren():
                                if c.tag == DRAWINGML('stCxn'):
                                    shape_properties['connection-start'] = self.__shape_id(c.attrib['id'])
                                elif c.tag == DRAWINGML('endCxn'):
                                    shape_properties['connection-end'] = self.__shape_id(c.attrib['id'])
                        shape_properties['line-style'] = pptx_shape.line.prstDash                   # type: ignore
                        shape_properties['head-end'] = pptx_shape.line.headEnd.get('type', 'none')  # type: ignore
                        shape_properties['tail-end'] = pptx_shape.line.tailEnd.get('type', 'none')  # type: ignore
                        shape_properties['stroke-width'] = abs(transform.scale_length((int(pptx_shape.line.width.emu), 0))[0])  # type: ignore
                        shape_properties['stroke-width'] /= STROKE_WIDTH_SCALE_FACTOR
                    else:
                        name = self.__text_content(pptx_shape)
                        shape_type = SHAPE_TYPE.COMPONENT
                        if name != '':
                            shape_properties['name'] = name
                            shape_properties['align'] = text_alignment(pptx_shape)
                    shape = self.__new_shape(pptx_shape.shape_id, geometry, shape_properties, shape_type)
                    shapes.append(shape)
                elif geometry is None:
                    log.warning(f'Shape "{shape_name}" {pptx_shape.shape_type}/{shape_properties.get("shape-kind")} not processed -- cannot get geometry')
                else:
                    log.warning(f'Shape "{shape_name}" {pptx_shape.shape_type}/{shape_properties.get("shape-kind")} not processed -- cannot get valid geometry')
            elif pptx_shape.shape_type == MSO_SHAPE_TYPE.GROUP:             # type: ignore
                shapes.append(self.__process_group(pptx_shape, transform))  # type: ignore
            elif pptx_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:           # type: ignore
                if good_geometry(geometry := get_shape_geometry(pptx_shape, transform, shape_properties)):
                    shape = self.__new_shape(pptx_shape.shape_id, geometry, shape_properties, SHAPE_TYPE.IMAGE)
                    bbox = geometry.bounds                      # type: ignore
                    image_pos = (bbox[0], bbox[1])
                    image_size = (bbox[2]-bbox[0], bbox[3]-bbox[1])
                    image = base64.b64encode(pptx_shape.image.blob).decode('utf-8')
                    image_data = f'data:{pptx_shape.image.content_type};charset=utf-8;base64,{image}'
                    image_rect = svgelements.Rect(*image_pos, *image_size)
                    image_rect.set('data-image-href', image_data)
                    shape.set_property('svg-element', image_rect)
                    shape.set_property('svg-kind', 'image')
                    shapes.append(shape)
            else:
                log.warning('Shape "{}" {} not processed...'.format(shape_name, str(pptx_shape.shape_type)))
            progress_bar.update(1)

        progress_bar.close()
        return shapes

#===============================================================================

class Powerpoint():
    def __init__(self, flatmap: 'FlatMap', source: 'PowerpointSource',      # type: ignore
                       SlideClass=Slide, slide_options: Optional[dict]=None):
        ppt_bytes = FilePath(source.href).get_BytesIO()
        pptx = Presentation(ppt_bytes)

        (width, height) = (pptx.slide_width, pptx.slide_height)
        self.__transform = Transform([[WORLD_METRES_PER_EMU,                     0, 0],
                                      [                    0, -WORLD_METRES_PER_EMU, 0],
                                      [                    0,                     0, 1]])@np.array([[1.0, 0.0,  -width/2.0],
                                                                                                    [0.0, 1.0, -height/2.0],
                                                                                                    [0.0, 0.0,         1.0]])
        top_left = self.__transform.transform_point((0, 0))
        bottom_right = self.__transform.transform_point((width, height))
        # southwest and northeast corners
        self.__bounds = (top_left[0], bottom_right[1], bottom_right[0], top_left[1])

        colour_theme = ColourTheme(ppt_bytes)
        if slide_options is None:
            slide_options = {}
        self.__slides: list[Slide] = [SlideClass(flatmap, source,
                                                 slide_index, slide,
                                                 colour_theme,
                                                 self.__bounds, self.__transform,
                                                 **slide_options)
                                            for slide_index, slide in enumerate(pptx.slides)]

    @property
    def bounds(self) -> MapBounds:
        return self.__bounds

    @property
    def slides(self) -> list[Slide]:
        return self.__slides

    @property
    def transform(self) -> Transform:
        return self.__transform

#===============================================================================
