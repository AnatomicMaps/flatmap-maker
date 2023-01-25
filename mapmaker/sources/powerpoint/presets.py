#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

import os.path

import pptx.oxml as oxml
import pptx.oxml.ns as ns

from pptx.dml.color import ColorFormat
from pptx.dml.line import LineFormat
from pptx.oxml.dml.color import CT_Percentage, _BaseColorElement
from pptx.oxml.theme import CT_OfficeStyleSheet

from pptx.oxml.shapes.autoshape import CT_GeomGuideList
from pptx.oxml.shapes.groupshape import CT_GroupShapeProperties
from pptx.oxml.shapes.shared import CT_LineProperties
from pptx.oxml.simpletypes import XsdString
from pptx.oxml.slide import _BaseSlideElement

from pptx.oxml.xmlchemy import (
    BaseOxmlElement,
    Choice,
    OneAndOnlyOne,
    RequiredAttribute,
    ZeroOrOne,
    ZeroOrOneChoice
)

#===============================================================================

PPTX_NAMESPACE = {
    'a': "http://schemas.openxmlformats.org/drawingml/2006/main",
    'p': "http://schemas.openxmlformats.org/presentationml/2006/main",
    'r': "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
}

def DRAWINGML(tag):
    return f'{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}'

def pptx_resolve(qname: str) -> str:
#===================================
    parts = qname.split(':', 1)
    if len(parts) == 2 and parts[0] in PPTX_NAMESPACE:
        return f'{{{PPTX_NAMESPACE[parts[0]]}}}{parts[1]}'
    return qname

def pptx_uri(qname: str) -> str:
#===============================
    parts = qname.split(':', 1)
    if len(parts) == 2 and parts[0] in PPTX_NAMESPACE:
        return f'{PPTX_NAMESPACE[parts[0]]}/{parts[1]}'
    return qname

#===============================================================================

ns._nsmap['drawml'] = ("http://www.ecma-international.org/flat/publications/standards/Ec"
                       "ma-376/drawingml/")

#===============================================================================

class PresetShapeDefinition(BaseOxmlElement):
    """`drawml:presetShapeDefinition` element class."""
    @classmethod
    def new(cls, xml):
        """Return shape definitions configured as ..."""
        return oxml.parse_xml(xml)

#===============================================================================

class Geometry2D(BaseOxmlElement):
    avLst = ZeroOrOne("a:avLst")
    gdLst = ZeroOrOne("a:gdLst")
    pathLst = ZeroOrOne("a:pathLst")

#===============================================================================

class PresetShape(Geometry2D):
    """`drawml:PresetShape` element class."""

    name = RequiredAttribute("name", XsdString)

#===============================================================================

oxml.register_element_cls("a:custGeom", Geometry2D)
oxml.register_element_cls("a:gdLst", CT_GeomGuideList)

oxml.register_element_cls("drawml:presetShapeDefinition", PresetShapeDefinition)
oxml.register_element_cls("drawml:presetShape", PresetShape)

#===============================================================================

class Shapes(object):
    definitions_ = {}

    with open(os.path.join(os.path.dirname(__file__), 'presetShapeDefinitions.xml'), 'rb') as defs:
        for defn in PresetShapeDefinition.new(defs.read()):
            definitions_[defn.name] = defn

    @staticmethod
    def lookup(name):
        return Shapes.definitions_[name]

#===============================================================================
#===============================================================================

class ThemeDefinition(CT_OfficeStyleSheet):
    name = RequiredAttribute("name", XsdString)
    themeElements = OneAndOnlyOne("a:themeElements")

    @classmethod
    def new(cls, xml):
        """Return theme definition"""
        t = oxml.parse_xml(xml)
        return t

#===============================================================================

class ThemeElements(BaseOxmlElement):
    clrScheme = OneAndOnlyOne("a:clrScheme")

#===============================================================================

class ColourScheme(BaseOxmlElement):
    name = RequiredAttribute("name", XsdString)

#===============================================================================

oxml.register_element_cls("a:theme", ThemeDefinition)
oxml.register_element_cls("a:themeElements", ThemeElements)
oxml.register_element_cls("a:clrScheme", ColourScheme)

#===============================================================================

class CT_SlideMasterUpdated(_BaseSlideElement):
    """
    ``<p:sldMaster>`` element, root of a slide master part
    """

    _tag_seq = (
        "p:cSld",
        "p:clrMap",
        "p:sldLayoutIdLst",
        "p:transition",
        "p:timing",
        "p:hf",
        "p:txStyles",
        "p:extLst",
    )
    cSld = OneAndOnlyOne("p:cSld")
    clrMap = OneAndOnlyOne("p:clrMap")    ### We need access to clrMap
    sldLayoutIdLst = ZeroOrOne("p:sldLayoutIdLst", successors=_tag_seq[3:])
    del _tag_seq

#===============================================================================

oxml.register_element_cls("p:sldMaster", CT_SlideMasterUpdated)

#===============================================================================

class CT_GroupShapePropertiesUpdated(CT_GroupShapeProperties):
    """p:grpSpPr element """

    _tag_seq = (
        "a:xfrm",
        "a:noFill",
        "a:solidFill",
        "a:gradFill",
        "a:blipFill",
        "a:pattFill",
        "a:grpFill",
        "a:effectLst",
        "a:effectDag",
        "a:scene3d",
        "a:extLst",
    )
    xfrm = ZeroOrOne("a:xfrm", successors=_tag_seq[1:])
    eg_groupFillProperties = ZeroOrOneChoice(
        (
            Choice("a:noFill"),
            Choice("a:solidFill"),
            Choice("a:gradFill"),
            Choice("a:blipFill"),
            Choice("a:pattFill"),
            Choice("a:grpFill"),
        ),
        successors=_tag_seq[7:],
    )
    effectLst = ZeroOrOne("a:effectLst", successors=_tag_seq[8:])
    del _tag_seq

    @property
    def eg_fillProperties(self):
        """
        Required to fulfill the interface used by dml.fill.
        """
        return self.eg_groupFillProperties

#===============================================================================

oxml.register_element_cls("p:grpSpPr", CT_GroupShapePropertiesUpdated)

#===============================================================================

# Monkey patching color to get colour properties...

oxml.register_element_cls('a:alpha', CT_Percentage)
_BaseColorElement.alpha = ZeroOrOne("a:alpha")
_BaseColorElement.alpha.populate_class_members(_BaseColorElement, "alpha")
ColorFormat.alpha = property(lambda self: (self._color._xClr.alpha.val
                                 if self._color._xClr.alpha is not None
                                 else 1.0))

ColorFormat.lumMod = property(lambda self: (self._color._xClr.lumMod.val
                                  if self._color._xClr.lumMod is not None
                                  else 1.0))
ColorFormat.lumOff = property(lambda self: (self._color._xClr.lumOff.val
                                  if self._color._xClr.lumOff is not None
                                  else 0.0))

oxml.register_element_cls("a:satMod", CT_Percentage)
_BaseColorElement.satMod = ZeroOrOne("a:satMod")
_BaseColorElement.satMod.populate_class_members(_BaseColorElement, "satMod")
ColorFormat.satMod = property(lambda self: (self._color._xClr.satMod.val
                                  if self._color._xClr.satMod is not None
                                  else 1.0))

oxml.register_element_cls("a:shade", CT_Percentage)
_BaseColorElement.shade = ZeroOrOne("a:shade")
_BaseColorElement.shade.populate_class_members(_BaseColorElement, "shade")
ColorFormat.shade = property(lambda self: (self._color._xClr.shade.val
                                  if self._color._xClr.shade is not None
                                  else 1.0))

oxml.register_element_cls("a:tint", CT_Percentage)
_BaseColorElement.tint = ZeroOrOne("a:tint")
_BaseColorElement.tint.populate_class_members(_BaseColorElement, "tint")
ColorFormat.tint = property(lambda self: (self._color._xClr.tint.val
                                  if self._color._xClr.tint is not None
                                  else 0.0))

#===============================================================================

# Monkey patching line properties to get end types...

CT_LineProperties.headEnd = ZeroOrOne("a:headEnd")
CT_LineProperties.headEnd.populate_class_members(CT_LineProperties, "headEnd")
LineFormat.headEnd = property(lambda self: (self._ln.headEnd.attrib
                                  if self._ln is not None and self._ln.headEnd is not None
                                  else {}))

CT_LineProperties.tailEnd = ZeroOrOne("a:tailEnd")
CT_LineProperties.tailEnd.populate_class_members(CT_LineProperties, "tailEnd")
LineFormat.tailEnd = property(lambda self: (self._ln.tailEnd.attrib
                                  if self._ln is not None and self._ln.tailEnd is not None
                                  else {}))

#===============================================================================
