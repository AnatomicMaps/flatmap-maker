#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019 - 2023  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from typing import Optional

#===============================================================================

from pptx.slide import Slide as PptxSlide

import shapely.geometry
import shapely.prepared
import shapely.strtree

#===============================================================================

from mapmaker.geometry import Transform
from mapmaker.knowledgebase.celldl import CD_CLASS, FC_CLASS, FC_KIND
from mapmaker.knowledgebase.sckan import SckanNeuronPopulations
from mapmaker.settings import settings
from mapmaker.sources.shape import Shape, SHAPE_TYPE
from mapmaker.sources import MapBounds
from mapmaker.utils import log, TreeList

from ..powerpoint import PowerpointSource, Slide
from ..shapefilter import ShapeFilter
from ..powerpoint.colour import ColourTheme

#===============================================================================

from .components import make_annotation, make_component, make_connection, make_connector
from .components import is_annotation, is_component, is_connector, is_system_name
from .components import ensure_parent_system
from .components import HYPERLINK_KINDS, HYPERLINK_IDENTIFIERS
from .components import NERVE_FEATURE_KINDS, NEURON_PATH_TYPES
from .components import ORGAN_COLOUR, ORGAN_KINDS
from .components import VASCULAR_KINDS, VASCULAR_REGION_COLOUR, VASCULAR_VESSEL_KINDS
from .connections import ConnectionClassifier

#===============================================================================

# Shapes smaller than this are assumed to be connectors or hyperlinks
MAX_CONNECTOR_AREA = 120000000      # metres**2

MIN_OVERLAP_FRACTION = 0.2  # Smaller geometry and >= 20% common area ==> containment
                            # e.g. `cochlear nuclei` and `pons`

MAX_FTU_OUTSIDE = 0.1       # An FTU should mainly be inside its organ

def contained_in(inside_shape, outer_shape, outside_fraction=0.0):
    if outside_fraction== 0.0:
        return outer_shape.geometry.contains(inside_shape.geometry)
    excess = inside_shape.geometry.difference(outer_shape.geometry)
    return excess.area <= outside_fraction*inside_shape.geometry.area

#===============================================================================

class FCPowerpointSource(PowerpointSource):
    def __init__(self, flatmap, id, href, kind, source_range=None,
                 shape_filter: Optional[ShapeFilter]=None, **kwds):
        super().__init__(flatmap, id, href, kind=kind,
                         source_range=source_range,
                         SlideClass=FCSlide, slide_options=dict(
                            shape_filter=shape_filter,
                            sckan_neurons=flatmap.sckan_neuron_populations
                         ),
                         **kwds)

    def get_raster_source(self):
    #===========================
        return None                 # We don't rasterise FC maps

#===============================================================================

# The outer geometry of the slide
SLIDE_LAYER_ID = 'SLIDE-LAYER-ID'

#===============================================================================

class FCSlide(Slide):
    def __init__(self, flatmap: 'FlatMap', source: PowerpointSource, index: int, pptx_slide: PptxSlide,   # type: ignore
                 theme: ColourTheme, bounds: MapBounds, transform: Transform,
                 shape_filter: Optional[ShapeFilter]=None,
                 sckan_neurons: Optional[SckanNeuronPopulations]=None):
        super().__init__(flatmap, source, index, pptx_slide, theme, bounds, transform)
        self.__shape_filter = shape_filter
        self.__sckan_neurons = sckan_neurons
        self.__shapes_by_id: dict[str, Shape] = {
            SLIDE_LAYER_ID: make_component(Shape(SHAPE_TYPE.LAYER, SLIDE_LAYER_ID, self.geometry,
                                            {'name': f'{source.id.capitalize()} Layer'}))
        }
        self.__connection_classifier = ConnectionClassifier()
        self.__connections = []
        self.__organ_ids: set[str] = set()
        self.__nerve_ids: set[str] = set()
        self.__system_ids: set[str] = set()

    def process(self, annotator: Optional['Annotator']=None):
    #======================================================
        super().process(annotator)

        self.__classify_shapes()
        if annotator is not None:
            self.__add_annotation(annotator)

        if self.__shape_filter is not None:
            if self.kind == 'base':
                # Add shapes to the filter
                for shape in self.shapes.flatten(skip=1):
                    self.__shape_filter.add_shape(shape)
                # Now create it for use by subsequent layers
                self.__shape_filter.create_filter()
            elif self.kind == 'layer':
                # Exclude shapes from the layer if they are similar to those in the base layer.
                # Excluded shapes have a ``global-shape`` property giving the matching base shape.
                for shape in self.shapes.flatten(skip=1):
                    self.__shape_filter.filter(shape)

        self.__add_connections()
        return self.shapes

    def __classify_shapes(self):
    #===========================
        # First extract shape geometries and create a spatial index
        # so we can find their containment hierarchy

        geometry_to_shape = {}
        geometries = []
        def add_shape_geometry(fc_shape):
            geometry_to_shape[id(geometry)] = fc_shape
            geometries.append(geometry)

        outer_geometry = shapely.prepared.prep(self.geometry)
        for shape in self.shapes.flatten(skip=1):
            geometry = shape.geometry
            if shape.type == SHAPE_TYPE.FEATURE and 'Polygon' in geometry.geom_type:
                # We are only interested in features actually on the slide that are
                # either components or connectors
                if outer_geometry.contains(geometry):
                    shape_kind = shape.properties.get('shape-kind', '')
                    if shape.colour is None:
                        if shape.name != '':
                            fc_shape = make_annotation(shape, FC_CLASS.DESCRIPTION)
                            fc_shape.set_property('exclude', True)    ## Only include if authoring??
                    elif (geometry.area < MAX_CONNECTOR_AREA):
                        fc_shape = None
                        if shape_kind.startswith('star'):
                            if (kind := HYPERLINK_KINDS.lookup(shape.colour)) is not None:
                                fc_shape = make_annotation(shape, FC_CLASS.HYPERLINK)
                                fc_shape.fc_kind = kind
                        else:
                            fc_shape = make_connector(shape)
                        if fc_shape is not None:
                            self.__shapes_by_id[shape.id] = fc_shape
                            add_shape_geometry(fc_shape)
                    else:
                        fc_shape = make_component(shape)
                        self.__shapes_by_id[shape.id] = fc_shape
                        add_shape_geometry(fc_shape)
            elif shape.type == SHAPE_TYPE.CONNECTION:
                self.__connections.append(make_connection(shape))

        # Spatial index to find component containment hierarchy
        idx = shapely.strtree.STRtree(geometries)

        # We now identify systems and for non-system features (both components and connectors)
        # find features which overlap them
        cardio_system = None
        nervous_system = None
        non_system_components = []
        connectors = []
        hyperlinks = []
        for shape_id, fc_shape in self.__shapes_by_id.items():
            # Do we need a better way of detecting systems??
            if is_component(fc_shape) and is_system_name(fc_shape.name):
                fc_shape.fc_class = FC_CLASS.SYSTEM
                fc_shape.add_parent(self.__shapes_by_id[SLIDE_LAYER_ID])
                self.__system_ids.add(shape_id)
                if 'CARDIO' in fc_shape.name:
                    fc_shape.fc_kind = FC_KIND.CARDIOVASCULAR_SYSTEM
                    cardio_system = fc_shape
                elif 'BRAIN' in fc_shape.name:
                    fc_shape.fc_kind = FC_KIND.NERVOUS_SYSTEM
                    nervous_system = fc_shape
            else:       # Component, Connector, or Annotation (Hyperlink)
                # STRtree query returns geometries whose bounding box intersects the shape's bounding box
                bigger_intersecting_geometries: list[int] = [id for id in idx.query(fc_shape.geometry)
                                                        if geometries[id].area > fc_shape.geometry.area
                                                        and geometries[id].intersection(fc_shape.geometry).area  # type: ignore
                                                            >= MIN_OVERLAP_FRACTION*fc_shape.geometry.area]
                # Set the shape's parents, ordered by the area of its overlapping geometries,
                # with the smallest (immediate) parent first
                containing_ids_area_order = [id_area[0]
                    for id_area in sorted([(id(geometries[index]), geometries[index].area)
                        for index in bigger_intersecting_geometries], key = lambda x: x[1])]
                if is_component(fc_shape):
                    if len(containing_ids_area_order):
                        parent = geometry_to_shape[containing_ids_area_order[0]]
                        if parent.fc_class != FC_CLASS.SYSTEM:      # Systems are only parents of Organs
                            fc_shape.add_parent(parent)             # and are assigned later
                        fc_shape.containing_ids = containing_ids_area_order
                    non_system_components.append(fc_shape)

                elif is_connector(fc_shape):
                    parent = None
                    for shape_id in containing_ids_area_order:
                        parent = geometry_to_shape[shape_id]
                        if is_component(parent):
                            break
                    if parent is not None:
                        fc_shape.add_parent(parent)
                        fc_shape.properties['parent-id'] = parent.id
                    else:
                        fc_shape.log_error(f'Connector has no parent: {fc_shape}')
                    connectors.append(fc_shape)
                elif is_annotation(fc_shape) and fc_shape.fc_class == FC_CLASS.HYPERLINK:
                    fc_shape.add_parent(geometry_to_shape[containing_ids_area_order[0]])
                    hyperlinks.append(fc_shape)

        # Classify connectors that are unambigously neural connectors
        for connector in connectors:
            if (connector.shape_kind == 'rect'
            and (neuron_type := NEURON_PATH_TYPES.lookup(connector.colour)) is not None):
                connector.fc_class = FC_CLASS.NEURAL
                connector.fc_kind = FC_KIND.CONNECTOR_PORT
                connector.path_type = neuron_type

        # Classify as FTUs those components that have a neural connector port
        # which has the component as their immediate parent
        for fc_shape in non_system_components:
            for child in fc_shape.children:
                if is_connector(child) and child.fc_kind == FC_KIND.CONNECTOR_PORT:
                    fc_shape.fc_class = FC_CLASS.FTU
                    break

        # Organs are components that are contained in at least one system and
        # have ORGAN_COLOUR colour or contain FTUs
        for fc_shape in non_system_components:
            if fc_shape.fc_class == FC_CLASS.UNKNOWN:
                have_organ = False
                organ_kind = FC_KIND.UNKNOWN
                if (ORGAN_COLOUR.matches(fc_shape.colour)
                 or (organ_kind := ORGAN_KINDS.lookup(fc_shape.colour)) is not None):
                    have_organ = True
                else:
                    for child in fc_shape.children:
                        if (child.fc_class == FC_CLASS.FTU
                        and contained_in(child, fc_shape, MAX_FTU_OUTSIDE)):
                            have_organ = True
                            break
                if have_organ:
                    fc_shape.fc_class = FC_CLASS.ORGAN
                    fc_shape.fc_kind = organ_kind
                    self.__organ_ids.add(fc_shape.id)
                    if fc_shape.name == '':
                        fc_shape.log_error(f'An organ must have a name: {fc_shape}')
                    have_system = False
                    for shape_id in fc_shape.containing_ids:
                        parent = geometry_to_shape[shape_id]
                        if parent.fc_class == FC_CLASS.SYSTEM:      # Systems are only parents of Organs
                            fc_shape.add_parent(parent)
                            have_system = True
                    if not have_system:
                        fc_shape.log_error(f'An organ must be in at least one system: {fc_shape}')

        # Components within an organ are either vascular regions or FTUs
        for fc_shape in non_system_components:
            if (fc_shape.fc_class == FC_CLASS.UNKNOWN
            and len(fc_shape.parents) and (parent := fc_shape.parents[0]).id in self.__organ_ids):
                if VASCULAR_REGION_COLOUR.matches(fc_shape.colour):
                    fc_shape.fc_class = FC_CLASS.VASCULAR
                    fc_shape.fc_kind = FC_KIND.VASCULAR_REGION
                else:
                    fc_shape.fc_class = FC_CLASS.FTU
            # Check vascular  regions and FTUs are only in a single organ
            if fc_shape.fc_class in [FC_CLASS.FTU, FC_CLASS.VASCULAR]:
                for parent in fc_shape.parents[1:]:
                    if parent.id in self.__organ_ids:
                        fc_shape.log_error(f'FTUs and regions can only be in a single organ: {fc_shape}')
                        break

        # Remaining named components should be either neural or vascular
        for fc_shape in non_system_components:
            if fc_shape.fc_class == FC_CLASS.UNKNOWN:
                if ((kind := NERVE_FEATURE_KINDS.lookup(fc_shape.colour)) is not None
                 or nervous_system in fc_shape.parents):
                    if fc_shape.name == 'G':
                        fc_shape.name = 'Ganglion'
                    fc_shape.fc_class = FC_CLASS.NEURAL
                    self.__nerve_ids.add(fc_shape.id)
                    fc_shape.description = kind
                elif (kind := VASCULAR_VESSEL_KINDS.lookup(fc_shape.colour)) is not None:
                    fc_shape.fc_class = FC_CLASS.VASCULAR
                    fc_shape.fc_kind = kind
            elif fc_shape.fc_class == FC_CLASS.UNKNOWN:
                unknown_shapes.append(fc_shape)
            # A nerve or vessel that connections may pass through
            if fc_shape.fc_class in [FC_CLASS.NEURAL, FC_CLASS.VASCULAR]:
                fc_shape.cd_class = CD_CLASS.CONDUIT
                self.__connection_classifier.add_component(fc_shape)
                # All neural components must be part of the nervous system and
                # all vascular components must be part of the cardiovascular system
                ensure_parent_system(fc_shape, nervous_system if fc_shape.fc_class == FC_CLASS.NEURAL
                                               else cardio_system)

        # Hyperlinks become properties of the feature they are on
        for hyperlink in hyperlinks:
            if hyperlink.has_property('hyperlink') and (parent := hyperlink.parent) is not None:
                parent.get_property('hyperlinks').append({
                    'id': HYPERLINK_IDENTIFIERS[hyperlink.fc_kind],
                    'url': hyperlink.properties['hyperlink']
                })
            hyperlink.set_property('exclude', True)

        # Classify remaining connectors
        for connector in connectors:
            if connector.fc_class == FC_CLASS.UNKNOWN:
                if connector.shape_kind == 'leftRightArrow':
                    connector.fc_class = FC_CLASS.NEURAL
                    connector.fc_kind = FC_KIND.CONNECTOR_JOINER
                elif connector.parent is not None and connector.parent.fc_class == FC_CLASS.NEURAL:
                    if connector.shape_kind == 'ellipse':
                        if (path_type := NEURON_PATH_TYPES.lookup(connector.colour)) is not None:
                            connector.fc_class = FC_CLASS.NEURAL
                            connector.fc_kind = FC_KIND.CONNECTOR_NODE
                            connector.path_type = path_type
                    elif connector.shape_kind == 'plus':
                        connector.fc_class = FC_CLASS.NEURAL
                        connector.fc_kind = FC_KIND.PLEXUS
                elif connector.shape_kind == 'ellipse':
                    if (kind := VASCULAR_KINDS.lookup(connector.colour)) is not None:
                        connector.fc_class = FC_CLASS.VASCULAR
                        connector.fc_kind = FC_KIND.CONNECTOR_NODE
                        connector.description = kind
                    elif (path_type := NEURON_PATH_TYPES.lookup(connector.colour)) is not None:
                        connector.fc_class = FC_CLASS.NEURAL
                        connector.path_type = path_type
                        if connector.parent is not None and connector.parent.fc_class == FC_CLASS.FTU:
                            connector.fc_kind = FC_KIND.GANGLION
                            connector.name = f'{path_type.name} ganglion'
                        else:
                            # Connector either without a parent or in a non-neural, non-FTU component
                            connector.shape.log_error(f'Connector not in FTU as expected: {connector.shape}')
                            connector.fc_kind = FC_KIND.CONNECTOR_NODE
            if connector.fc_class != FC_CLASS.UNKNOWN:
                self.__connection_classifier.add_connector(connector)

    def __add_annotation(self, annotator: 'Annotator'):
    #================================================
        # Called after shapes have been extracted
        for fc_shape in self.__shapes_by_id.values():
            if (is_component(fc_shape)
            and (term := annotator.lookup_shape(fc_shape)) is not None):
                fc_shape.properties['models'] = term

        # go through all connectors and set FTU/organ for them
        for fc_shape in self.__shapes_by_id.values():
            if is_connector(fc_shape):
                if (fc_shape.fc_kind == FC_KIND.GANGLION
                and (term := annotator.lookup_shape(fc_shape)) is not None):
                    fc_shape.properties['models'] = term
                self.__annotate_connector(fc_shape)

    def __annotate_connector(self, connector: Shape):
    #================================================
        names = [connector.name.capitalize()] if connector.name else []
        parent_models = []
        def set_label(parent):
            if parent.models:
                parent_models.append(parent.models)
            names.append(parent.name.capitalize())

        if connector.parent is not None:
            set_label(connector.parent)
            if connector.parent.parents:
                set_label(connector.parent.parents[0])
            connector.properties['name'] = '/'.join(names)
            if len(parent_models):
                connector.properties['parent-models'] = tuple(parent_models)
        connector.shape_type = 'connector'

    def __feature_properties(self, feature_id):
    #==========================================
        if (feature := self.flatmap.get_feature(feature_id)) is not None:
            return feature.properties
        elif (shape := self.__shapes_by_id.get(feature_id)) is not None:
            return shape.properties
        return {}

    def __add_connections(self):
    #===========================
        for connection in self.__connections:
            # First pass will join together sub-paths at joining nodes
            self.__connection_classifier.add_connection(connection)

        for connection in self.__connections:
            end_names = []
            end_node_parents = set()
            for connector_id in connection.connector_ids:
                properties = self.__feature_properties(connector_id)
                if (parent_id := properties.get('parent-id')) is not None:
                    end_node_parents.add(parent_id)
                if settings.get('authoring', False) and (name := properties.get('name', '')):
                    end_names.append(f'CN: {name.capitalize()}')
            connection.properties['name'] = '\n'.join(end_names)
            connection.properties['node-ids'] = list(end_node_parents
                                                   | set(connection.connector_ids)
                                                   | set(connection.intermediate_components)
                                                   | set(connection.intermediate_connectors))

            if len([connector for connector in connection.get_property('connectors', [])
                if connector.fc_kind in [FC_KIND.CONNECTOR_FREE_END, FC_KIND.CONNECTOR_JOINER]]):
                    log.warning(f'Connection has free ends: {connection}')
                    connection.properties['kind'] = 'error'

            # Save neuron paths for generating connectivity
            if self.__sckan_neurons is not None and connection.fc_kind == FC_KIND.NEURON:
                connection.properties['sckan'] = False                      # Assume no paths are valid
                self.__sckan_neurons.add_connection(self.__feature_properties, connection)

#===============================================================================
