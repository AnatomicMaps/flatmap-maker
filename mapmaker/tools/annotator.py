#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from collections import OrderedDict

#===============================================================================

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

#===============================================================================

class Slide(object):
    def __init__(self, slide):
        self._slide_id = slide.slide_id
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            self._notes = notes_slide.notes_text_frame.text.rstrip()
        else:
            self._notes = ''
        self._shapes_by_id = OrderedDict()
        self.process_shape_list_(slide.shapes)

    @property
    def notes(self):
        return self._notes

    @property
    def shapes_by_id(self):
        return self._shapes_by_id

    def process_shape_list_(self, shapes):
        for shape in shapes:
            unique_id = '{}#{}'.format(self._slide_id, shape.shape_id)
            self._shapes_by_id[unique_id] = shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.process_shape_list_(shape.shapes)

    def list(self):
        print('SLIDE: {!s:8} {}'.format(self._slide_id, self._notes))
        for id, shape in self._shapes_by_id.items():
            if (shape.name.startswith('.')
             or shape.name.startswith('#')):
                print('SHAPE: {!s:8} {}'.format(id, shape.name))

#===============================================================================

class Presentation(object):
    def __init__(self, powerpoint):
        self._pptx = pptx.Presentation(powerpoint)
        self._slides_by_id = OrderedDict()
        self._seen = []
        for slide in self._pptx.slides:
            self._slides_by_id[slide.slide_id] = Slide(slide)

    def list(self):
        for slide in self._slides_by_id.values():
            slide.list()

#===============================================================================

if __name__ == '__main__':
    import argparse
    import os, sys

    parser = argparse.ArgumentParser(description='Modify annotations in Powerpoint slides.')

    parser.add_argument('-l', '--list', action='store_true',
                        help='list annotations')

    parser.add_argument('-r', '--replace', nargs=2, action='append',
                        help='find and replace text using ')

    parser.add_argument('powerpoint', help='Powerpoint file')

    args = parser.parse_args()

    presentation = Presentation(args.powerpoint)

    if args.list:
        presentation.list()
