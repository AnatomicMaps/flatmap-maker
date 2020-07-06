#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from math import sqrt, sin, cos, pi as PI
import os

#===============================================================================

import numpy as np

import pptx.shapes.connector
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tqdm import tqdm

#===============================================================================

from parser import Parser
from properties import Properties

#===============================================================================

# Internal PPT units are EMUs (English Metric Units)

EMU_PER_CM  = 360000
EMU_PER_IN  = 914400

DOTS_PER_IN = 96

EMU_PER_DOT = EMU_PER_IN/DOTS_PER_IN

#===============================================================================

def cm_coords(x, y):
#===================
    return (x/EMU_PER_CM, y/EMU_PER_CM)

def ellipse_point(a, b, theta):
#==============================
    a_sin_theta = a*sin(theta)
    b_cos_theta = b*cos(theta)
    circle_radius = sqrt(a_sin_theta**2 + b_cos_theta**2)
    return (a*b_cos_theta/circle_radius, b*a_sin_theta/circle_radius)

#===============================================================================

class Transform(object):
    def __init__(self, shape, bbox=None):
        xfrm = shape.element.xfrm

        # From Section L.4.7.6 of ECMA-376 Part 1
        (Bx, By) = ((xfrm.chOff.x, xfrm.chOff.y)
                        if xfrm.chOff is not None else
                    (0, 0))
        (Dx, Dy) = ((xfrm.chExt.cx, xfrm.chExt.cy)
                        if xfrm.chExt is not None else
                    bbox)
        (Bx_, By_) = (xfrm.off.x, xfrm.off.y)
        (Dx_, Dy_) = (xfrm.ext.cx, xfrm.ext.cy)
        theta = xfrm.rot*PI/180.0
        Fx = -1 if xfrm.flipH else 1
        Fy = -1 if xfrm.flipV else 1
        T_st = np.array([[Dx_/Dx,      0, Bx_ - (Dx_/Dx)*Bx] if Dx != 0 else [1, 0, Bx_],
                         [     0, Dy_/Dy, By_ - (Dy_/Dy)*By] if Dy != 0 else [0, 1, By_],
                         [     0,      0,                 1]])
        U = np.array([[1, 0, -(Bx_ + Dx_/2.0)],
                      [0, 1, -(By_ + Dy_/2.0)],
                      [0, 0,                1]])
        R = np.array([[cos(theta), -sin(theta), 0],
                      [sin(theta),  cos(theta), 0],
                      [0,                    0, 1]])
        Flip = np.array([[Fx,  0, 0],
                         [ 0, Fy, 0],
                         [ 0,  0, 1]])
        T_rf = np.linalg.inv(U)@R@Flip@U
        self.__T = T_rf@T_st

    def matrix(self):
        return self.__T

#===============================================================================

class Feature(object):
    def __init__(self, id, geometry, properties, has_children=False):
        self.__id = id
        self.__geometry = geometry
        self.__properties = properties.copy()
        self.__properties['id'] = id
        self.__has_children = has_children

    def __str__(self):
        return '{}: {}'.format(self.__geometry.geom_type, self.__properties)

    @property
    def annotated(self):
        return self.shape_name.startswith('.')

    @property
    def geometry(self):
        return self.__geometry

    @geometry.setter
    def geometry(self, geometry):
        self.__geometry = geometry

    @property
    def geom_type(self):
        return self.__geometry.geom_type if self.__geometry else None

    @property
    def has_children(self):
        return self.__has_children

    @property
    def id(self):
        return self.__id

    @property
    def feature_id(self):
        return self.__id.split('#')[-1]

    @property
    def shape_name(self):
        return self.__properties.get('shape_name', '')

    @property
    def properties(self):
        return self.__properties

    def has(self, property):
        return self.__properties.get(property, '') != ''

    def is_a(self, property):
        return self.__properties.get(property, False)

    def property(self, property, default=None):
        return self.__properties.get(property, default)

#===============================================================================

class FeaturesValueError(ValueError):
    def __init__(self, msg, features):
        super().__init__('\n  '.join([msg] + [str(f) for f in features]))

#===============================================================================

class Layer(object):
    def __init__(self, extractor, slide, slide_number):
        self._slide = slide
        self.__extractor = extractor
        self.__slide_number = slide_number
        self.__errors = []
        self.__external_properties = Properties(extractor.settings)
        self.__pathways = self.__external_properties.pathways

        # Find `layer-id` text boxes so we have a valid ID **before** using
        # it when setting a shape's `path_id`.
        self.__set_defaults()
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text
            if notes_text.startswith('.'):
                layer_directive = Parser.layer_directive(notes_text)
                if 'error' in layer_directive:
                    self.__errors.append('Slide {}: invalid layer directive: {}'
                                        .format(slide_number, notes_text))
                else:
                    self.__layer_id = layer_directive.get('id')
                self.__description = layer_directive.get('description', self.__layer_id.capitalize())
                self.__models = layer_directive.get('models', '')
                self.__background_for = layer_directive.get('background-for', '')
                self.__selectable = self.__background_for == '' and not layer_directive.get('not-selectable')
                self.__selected = layer_directive.get('selected', False)
                self.__queryable_nodes = layer_directive.get('queryable-nodes', False)
                self.__zoom = layer_directive.get('zoom', None)

        self.__map_features = []
#*        self.__ontology_data = self.settings.ontology_data
        self.__annotations = {}
        self.__current_group = []

    def __set_defaults(self):
    #========================
        self.__layer_id = 'layer-{:02d}'.format(self.__slide_number)
        self.__description = 'Layer {}'.format(self.__slide_number)
        self.__models = ''
        self.__background_for = ''
        self.__selectable = True
        self.__selected = False
        self.__queryable_nodes = False
        self.__zoom = None
        # Cannot overlap with slide shape ids...
        self.__next_local_id = 100001

    def __set_feature_id(self, feature):
    #===================================
        if feature.has('external-id'):
            self.__external_properties.set_feature_id(feature.property('external-id'), feature.id)
        if feature.has('class'):
            self.__external_properties.set_class_id(feature.property('class'), feature.id)

    @property
    def extractor(self):
        return self.__extractor

    @property
    def settings(self):
        return self.__extractor.settings

    @property
    def annotations(self):
        return self.__annotations

    @property
    def description(self):
        return self.__description

    @property
    def models(self):
        return self.__models

    @property
    def background_for(self):
        return self.__background_for

    @property
    def selected(self):
        return self.__selected

    @property
    def selectable(self):
        return self.__selectable

    @property
    def queryable_nodes(self):
        return self.__queryable_nodes

    @property
    def zoom(self):
        return self.__zoom

    @property
    def errors(self):
        return self.__errors

    @property
    def map_features(self):
        return self.__map_features

    @property
    def layer_id(self):
        return self.__layer_id

    @property
    def slide_id(self):
        return self._slide.slide_id

    @property
    def resolved_pathways(self):
        return self.__pathways.resolved_pathways if self.__pathways is not None else None

    def unique_id(self, id):
    #=======================
        return '{}#{}'.format(self.slide_id, id)

    def next_local_id(self):
    #=======================
        id = self.unique_id(self.__next_local_id)
        self.__next_local_id += 1
        return id

    def process_initialise(self):
    #============================
        self.__current_group.append('SLIDE')

    def process(self):
    #=================
        self.process_initialise()
        self.process_shape_list(self._slide.shapes, outermost=True)
        self.process_finialise()

    def process_finialise(self):
    #===========================
        self.__external_properties.set_feature_ids()

    def save(self, filename=None):
    #=============================
        # Override in sub-class
        pass

    def process_group(self, group, properties, *args):
    #=================================================
        self.process_shape_list(group.shapes, *args)

    def process_shape(self, shape, properties, *args):
    #=================================================
        # Override in sub-class
        pass

    def process_shape_list(self, shapes, *args, outermost=False):
    #============================================================
        if not self.__selectable:
            return []

        if outermost:
            print('Processing shape list...')
            progress_bar = tqdm(total=len(shapes),
                unit='shp', ncols=40,
                bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')

        features = []
        for shape in shapes:
            properties = self.__external_properties.get_properties(shape,
                            self.__current_group[-1],
                            self.__slide_number)
            if 'path' in properties:
                path_id = properties['path']
                self.__pathways.add_path(path_id)
            elif (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or isinstance(shape, pptx.shapes.connector.Connector)):
                geometry = self.process_shape(shape, properties, *args)
                feature = Feature(self.unique_id(shape.shape_id), geometry, properties)
                self.__set_feature_id(feature)
                features.append(feature)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.__current_group.append(properties.get('shape_name', "''"))
                grouped_feature = self.process_group(shape, properties, *args)
                self.__current_group.pop()
                if grouped_feature is not None:
                    self.__set_feature_id(grouped_feature)
                    features.append(grouped_feature)
            elif (shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
               or shape.shape_type == MSO_SHAPE_TYPE.PICTURE):
                pass
            else:
                print('"{}" {} not processed...'.format(shape.name, str(shape.shape_type)))
            if outermost:
                progress_bar.update(1)

        if outermost:
            progress_bar.close()
        return features

#===============================================================================

class Extractor(object):
    def __init__(self, pptx, settings, layer_class=Layer):
        self.__LayerClass = layer_class
        self.__pptx = Presentation(pptx)
        self.__settings = settings
        self.__slides = self.__pptx.slides
        self.__slide_size = [self.__pptx.slide_width, self.__pptx.slide_height]
        self.__layers = {}

    def __len__(self):
        return len(self.__slides)

    @property
    def layers(self):
        return self.__layers

    @property
    def settings(self):
        return self.__settings

    @property
    def slide_maker(self):
        return self.__slide_maker

    @property
    def slide_size(self):
        return self.__slide_size

    def bounds(self):
        return (0, 0, self.__slide_size[0], self.__slide_size[1])

    def slide(self, slide_number):
        return self.__slides[slide_number - 1]

    def slide_to_layer(self, slide_number, debug_xml=False):
        slide = self.slide(slide_number)
        if debug_xml:
            xml = open(os.path.join(self.__settings.output_dir, 'layer{:02d}.xml'.format(slide_number)), 'w')
            xml.write(slide.element.xml)
            xml.close()
        if self.__LayerClass is not None:
            layer = self.__LayerClass(self, slide, slide_number)
            layer.process()
            print('Slide {}, layer {}'.format(slide_number, layer.layer_id))

            if layer.layer_id in self.__layers:
                raise KeyError('Duplicate layer id ({}) in slide {}'.format(layer.layer_id, slide_number))
            self.__layers[layer.layer_id] = layer
            return layer

    def slides_to_layers(self, slide_range):
        if slide_range is None:
            slide_range = range(1, len(self.__slides)+1)
        elif isinstance(slide_range, int):
            slide_range = [slide_range]
        for n in slide_range:
            self.slide_to_layer(n)

#===============================================================================
