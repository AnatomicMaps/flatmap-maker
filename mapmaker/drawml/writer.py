#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

#===============================================================================

class SlideShapes(object):
    def __init__(self, slide):
        self._slide_id = slide.slide_id
        self._shapes_by_id = {}
        self.process_shape_list_(slide.shapes)

    @property
    def shapes_by_id(self):
        return self._shapes_by_id

    def process_shape_list_(self, shapes):
        for shape in shapes:
            unique_id = '{}-{}'.format(self._slide_id, shape.shape_id)
            self._shapes_by_id[unique_id] = shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.process_shape_list_(shape.shapes)

#===============================================================================

class AnnotationWriter(object):
    def __init__(self, pptx):
        self._pptx = Presentation(pptx)
        self._shapes_by_id = {}
        self._seen = []
        for slide in self._pptx.slides:
            slide_shapes = SlideShapes(slide)
            self._shapes_by_id.update(slide_shapes.shapes_by_id)

    def remove_unseen(self):
        for id, shape in self._shapes_by_id.items():
            if id not in self._seen and shape.name.startswith('#'):
                shape.name = shape.name[1:]

    def save(self, output_pptx):
        self._pptx.save(output_pptx)

    def update_annotation(self, id, text):
        shape = self._shapes_by_id.get(id)
        self._seen.append(id)
        if shape is not None and shape.name != text:
            shape.name = text

#===============================================================================
