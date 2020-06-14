#===============================================================================
#
#  Flatmap viewer and annotation tools
#
#  Copyright (c) 2019  David Brooks
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
#===============================================================================

import io
import json
from copy import deepcopy
from zipfile import ZipFile

#===============================================================================

import pptx

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tqdm import tqdm

#===============================================================================

from properties import Properties

#===============================================================================

EXCLUDE_SHAPE_TYPES = ['group', 'invisible', 'marker', 'path', 'region']
EXCLUDE_TILE_LAYERS = ['pathways']

#===============================================================================

LAYOUT_BLANK_SLIDE = 6

#===============================================================================

XPATH_GROUP_SPPR        = './p:grpSpPr'

XPATH_AUTOSHAPE_BY_ID_GROUP  = './p:sp/p:nvSpPr/p:cNvPr[@id={}]'
XPATH_AUTOSHAPE_BY_ID_SHAPE  = './p:nvSpPr/p:cNvPr[@id={}]'

XPATH_PICTURE_BY_ID_GROUP    = './p:pic/p:nvPicPr/p:cNvPr[@id={}]'
XPATH_PICTURE_BY_ID_SHAPE    = './p:nvPicPr/p:cNvPr[@id={}]'

XPATH_CONNECTION_BY_ID_GROUP = './p:cxnSp/p:nvCxnSpPr/p:cNvPr[@id={}]'
XPATH_CONNECTION_BY_ID_SHAPE = './p:nvCxnSpPr/p:cNvPr[@id={}]'

XPATH_CONNECTION_TYPE        = './p:spPr/a:prstGeom[@prst]'

XPATH_SLIDE_cSld        = './p:cSld'
XPATH_SLIDE_extLst      = './p:cSld/p:extLst'

XPATH_PRS_extLst        = './p:extLst'

#===============================================================================

class NameChecker(object):
    def __init__(self, properties):
        self.__properties = properties

    def valid(self, shape):
        for key, value in self.__properties.get_properties(shape).items():
            if key in EXCLUDE_SHAPE_TYPES:
                return False
            elif key == 'tile-layer' and value in EXCLUDE_TILE_LAYERS:
                return False
        return True

#===============================================================================

def valid_notes(notes):
    return notes and notes[0] in ['.', '#']

#===============================================================================

def clean_markup(text):
    return ('.' + text[1:]) if (text and text[0] == '#') else text

#===============================================================================

def connector_type(name):
    return (MSO_CONNECTOR.STRAIGHT if name == 'line'
       else MSO_CONNECTOR.ELBOW    if name == 'bentConnector3'
       else MSO_CONNECTOR.CURVE    if name == 'curvedConnector3'
       else MSO_CONNECTOR.MIXED)

#===============================================================================

class Presentation(object):
    def __init__(self, source_file, properties):
        print('Opening presentation...')
        self._source_file = source_file
        self._source = pptx.Presentation(source_file)
        self._name_checker = NameChecker(properties)

#        xml = open('dirty_prs.xml', 'w')
#        xml.write(self._source.element.xml)
#        xml.close()

        self._prs = pptx.Presentation()
        self._prs.slide_width = self._source.slide_width
        self._prs.slide_height = self._source.slide_height
        ##self._prs.slide_masters = self._source.slide_masters
        self._blank_layout = self._prs.slide_layouts[LAYOUT_BLANK_SLIDE]
        self._current_group = None
        self._group_stack = []
        self._clean_slide = None
        self._progress_bar = None

    def clean(self, output_file):
        for slide in self._source.slides:
            self.add_slide(slide)
        self.save(output_file)

    def save(self, output_file):
        if len(self._group_stack):
            raise ValueError('Unclosed group...')

        # Add `p:extLst` to new presentation
        self._prs.element.append(self._source.element.xpath(XPATH_PRS_extLst)[0])

#        xml = open('clean_prs.xml', 'w')
#        xml.write(self._prs.element.xml)
#        xml.close()

        output = io.BytesIO()
        self._prs.save(output)
        output.seek(0)

        print('Copying themes...')
        # First copy the contents of the saved presentation, apart
        # from its themes, to create a new PPTX file
        with ZipFile(output_file, 'w') as clean_pptx:
            with ZipFile(output, 'r') as saved_prs:
                for info in saved_prs.infolist():
                    if not info.filename.startswith('ppt/theme/'):
                        clean_pptx.writestr(info, saved_prs.read(info))

            # Finally copy the original themes to the new presentation
            with ZipFile(self._source_file, 'r') as source_prs:
                for info in source_prs.infolist():
                    if info.filename.startswith('ppt/theme/'):
                        clean_pptx.writestr(info, source_prs.read(info))

    def add_slide(self, slide):
#        xml = open('dirty_slide.xml', 'w')
#        xml.write(slide.element.xml)
#        xml.close()
        print('Cleaning slide...')
        self._clean_slide = self._prs.slides.add_slide(self._blank_layout)
        self._clean_slide.shapes.turbo_add_enabled = True
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if valid_notes(notes):
                clean_notes_slide = self._clean_slide.notes_slide
                clean_notes_slide.notes_text_frame.text = clean_markup(notes)

        self.start_shapes_(len(slide.shapes))
        self.add_shapes_(slide.shapes, slide=True)
        self.end_shapes_(slide)

        # Add `p:extLst` to new `cSld` element  ### ????
        new_csld_element = self._clean_slide.element.xpath(XPATH_SLIDE_cSld)[0]
        new_csld_element.append(slide.element.xpath(XPATH_SLIDE_extLst)[0])

#        xml = open('clean_slide.xml', 'w')
#        xml.write(self._clean_slide.element.xml)
#        xml.close()

    def start_shapes_(self, total_shapes):
        print('Cleaning shapes...')
        self._progress_bar = tqdm(total=total_shapes,
            unit='shp', ncols=40,
            bar_format='{l_bar}{bar}| {n_fmt}/{total_fmt}')
        self._current_group = self._clean_slide
        self._group_stack = []

    def end_shapes_(self, slide):
        # replace <p:grpSpPr> element of the slide's shapes group
        clean_shapes_element = self._clean_slide.shapes.element
        clean_shapes_element.replace(clean_shapes_element.xpath(XPATH_GROUP_SPPR)[0],
                                     slide.shapes.element.xpath(XPATH_GROUP_SPPR)[0])
        self._progress_bar.close()

    def add_shapes_(self, shapes, slide=False):
        for shape in shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM
             or shape.shape_type == MSO_SHAPE_TYPE.PICTURE
             or isinstance(shape, pptx.shapes.connector.Connector)):
                if self._name_checker.valid(shape):
                    self.append_shape_(shape)

            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.start_group_()
                self.add_shapes_(shape.shapes)
                self.end_group_(shape)

            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                pass

            else:
                print('Unknown shape type {}, "{}"...'.format(str(shape.shape_type), shape.name))
            if slide:
                self._progress_bar.update(1)

    def append_shape_(self, shape):
        # We add a new shape and then replace
        # its xml element with the shape being appended
        """
$ python mapmaker/tools/cleaner.py --properties map_sources/rat/rat_flatmap_properties.json  \
                                  tests/sources/paths.pptx tests/sources/paths_cleaned.pptx

File "mapmaker/tools/cleaner.py", line 230, in append_shape_
    shape.end_x, shape.end_y)
  File "/Users/dave/.local/share/virtualenvs/map-maker-uZb0MJDL/lib/python3.7/site-packages/pptx/shapes/shapetree.py", line 263, in add_connector
    cxnSp = self._add_cxnSp(connector_type, begin_x, begin_y, end_x, end_y)
  File "/Users/dave/.local/share/virtualenvs/map-maker-uZb0MJDL/lib/python3.7/site-packages/pptx/shapes/shapetree.py", line 384, in _add_cxnSp
    id_, name, connector_type, x, y, cx, cy, flipH, flipV
  File "/Users/dave/.local/share/virtualenvs/map-maker-uZb0MJDL/lib/python3.7/site-packages/pptx/oxml/shapes/groupshape.py", line 51, in add_cxnSp
    prst = MSO_CONNECTOR_TYPE.to_xml(type_member)
  File "/Users/dave/.local/share/virtualenvs/map-maker-uZb0MJDL/lib/python3.7/site-packages/pptx/enum/base.py", line 205, in to_xml
    cls.validate(enum_val)
  File "/Users/dave/.local/share/virtualenvs/map-maker-uZb0MJDL/lib/python3.7/site-packages/pptx/enum/base.py", line 176, in validate
    "%s not a member of %s enumeration" % (value, cls.__name__)
ValueError: MIXED (-2) not a member of MSO_CONNECTOR_TYPE enumeration
        """
        if isinstance(shape, pptx.shapes.connector.Connector):
            cxn_name = shape.element.xpath(XPATH_CONNECTION_TYPE)[0].get('prst')
            new_shape = self._current_group.shapes.add_connector(connector_type(cxn_name),
                                                                 shape.begin_x, shape.begin_y,
                                                                 shape.end_x, shape.end_y)
            new_id_element_xpath = XPATH_CONNECTION_BY_ID_GROUP.format(new_shape.shape_id)
            old_id_element_xpath = XPATH_CONNECTION_BY_ID_SHAPE.format(shape.shape_id)
        elif isinstance(shape, pptx.shapes.picture.Picture):
            new_shape = self._current_group.shapes.add_picture(io.BytesIO(shape.image.blob),
                                                               shape.left, shape.top,
                                                               shape.width, shape.height)
            new_id_element_xpath = XPATH_PICTURE_BY_ID_GROUP.format(new_shape.shape_id)
            old_id_element_xpath = XPATH_PICTURE_BY_ID_SHAPE.format(shape.shape_id)
        else:
            new_shape = self._current_group.shapes.add_shape(shape.shape_type,
                                                             shape.left, shape.top,
                                                             shape.width, shape.height)
            new_id_element_xpath = XPATH_AUTOSHAPE_BY_ID_GROUP.format(new_shape.shape_id)
            if isinstance(shape, pptx.shapes.autoshape.Shape):
                old_id_element_xpath = XPATH_AUTOSHAPE_BY_ID_SHAPE.format(shape.shape_id)
            else:
                print('Unexpected shape type: {}'.format(type(shape)))
                import pdb; pdb.set_trace()

        new_id_element = self._current_group.shapes.element.xpath(new_id_element_xpath)[0]
        new_element = new_id_element.getparent().getparent()

        shape_element = deepcopy(shape.element)
        shape_id_element = shape_element.xpath(old_id_element_xpath)[0]
        shape_id_element.set('id', str(new_shape.shape_id))

        #print('')
        #print('Parent:', new_element.xml)
        #print(' Shape:', shape_element.xml)

        self._current_group.shapes.element.replace(
            new_element,
            shape_element
            )

    def start_group_(self):
        self._group_stack.append(self._current_group)
        self._current_group = self._current_group.shapes.add_group_shape()
        self._current_group.shapes.turbo_add_enabled = True

    def end_group_(self, group):
        self._current_group.name = group.name
        # replace <p:grpSpPr> element of group
        clean_shapes_element = self._current_group.shapes.element
        clean_shapes_element.replace(clean_shapes_element.xpath(XPATH_GROUP_SPPR)[0],
                                     group.shapes.element.xpath(XPATH_GROUP_SPPR)[0])
        self._current_group = self._group_stack.pop()

#===============================================================================

def clean_presentation(source, target, properties):
#==================================================
    presentation = Presentation(source, properties)
    presentation.clean(target)

#===============================================================================

if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description='Clean Powerpoint slides for generating flatmap image tiles.')

    parser.add_argument('--anatomical-map',
                        help='Excel spreadsheet file for mapping shape classes to anatomical entities')
    parser.add_argument('--properties',
                    help='JSON file specifying pathways')

    parser.add_argument('source_ppt', help='Powerpoint file to clean')
    parser.add_argument('cleaned_ppt', help='Cleaned Powerpoint to create')
    args = parser.parse_args()

    ## Option to remove paths??
    ## Then properties file only when removing paths
    args = parser.parse_args()
    args.label_database = 'labels.sqlite'
    external_properties = Properties(args)

    clean_presentation(args.source_ppt, args.cleaned_ppt, external_properties)

#===============================================================================
